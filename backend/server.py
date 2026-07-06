from fastapi import FastAPI, APIRouter, HTTPException, Request, Response, Depends, Cookie, Header, WebSocket, WebSocketDisconnect, UploadFile, File, Form
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import json
import logging
import uuid
import re
import random
import asyncio
import base64
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional, Any
from datetime import datetime, timezone, timedelta
import httpx

# Must run before importing ocr_service (and any other module that reads
# os.environ at import time) — otherwise GOOGLE_CLOUD_VISION_API_KEY / etc.
# from backend/.env are invisible to those modules even though they're set.
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

from email_service import send_consent_email, send_progress_digest
from firebase_auth import verify_firebase_id_token, firebase_configured
from classroom_live import classroom_manager
from ocr_service import (
    decode_image_base64,
    get_quota_status,
    safe_get_quota_status,
    run_hybrid_ocr,
    QuotaExceededError,
    RateLimitError,
    VISION_ENABLED,
    get_device_profile,
)
from key_concepts import extract_key_concepts, concepts_to_knowledge_units, PRESENTATION_META
from behavioral_signals import record_behavioral_signal, summarize_for_parent, CONSENT_TYPE as BEHAVIORAL_CONSENT_TYPE
from consent_gate import assert_can_collect_learning_data, get_data_consent_status
from teacher_billing import register_teacher_billing_routes, register_teacher_seat, get_teacher_plan
from iap_service import register_iap_routes
from push_service import send_fcm_to_owner, fcm_configured
from game_library import detect_language
from rewards import (
    COINS_PER_CORRECT,
    BADGE_CATALOG,
    check_and_award_badges,
    get_student_badges,
)
from journey_engine import (
    READING_DICTATION_STEPS,
    RECITATION_DICTATION_STEPS,
    TRACK_STEP_TABLES,
    generate_step_game,
    journey_status,
    step_lock_status,
    score_meets_step_threshold,
    next_step_after_pass,
    assign_bundle_indices,
    units_for_bundle,
    bundle_count_for_units,
    step_back_on_fail,
    max_step_for_track,
    get_step_table,
    parse_step_completed_at,
    _step_was_completed,
)

import google.genai as genai
from google.genai import types as genai_types

# Hard cap on any single LLM call. Above this, we short-circuit to the deterministic
# fallback so the request never blocks the flow (Cloudflare/edge idle timeouts, plus
# the fact that users perceive >~10s of "PREPARING…" as broken).
LLM_TIMEOUT_S = 9.0

_genai_client: Optional["genai.Client"] = None


def _get_genai_client() -> "genai.Client":
    """Lazily build a single reusable google-genai client. Reads GEMINI_API_KEY from
    the environment automatically (see backend/.env) — no more Emergent-issued key."""
    global _genai_client
    if _genai_client is None:
        _genai_client = genai.Client()
    return _genai_client


async def _gemini_generate(system_message: str, prompt: str, timeout: float = LLM_TIMEOUT_S) -> str:
    """Direct replacement for the old emergentintegrations LlmChat().send_message()
    call. Same shape (system prompt + one user turn -> plain text back), same
    9s hard timeout, so every call site below only needed a one-line swap."""
    client = _get_genai_client()
    response = await asyncio.wait_for(
        client.aio.models.generate_content(
            model="gemini-3-flash-preview",
            contents=prompt,
            config=genai_types.GenerateContentConfig(system_instruction=system_message),
        ),
        timeout=timeout,
    )
    return response.text


mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url, serverSelectionTimeoutMS=10000)
db = client[os.environ['DB_NAME']]

GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
GEMINI_ENABLED = bool(GEMINI_API_KEY) and not str(GEMINI_API_KEY).startswith("PASTE")

app = FastAPI()
api_router = APIRouter(prefix="/api")


# ============== MODELS ==============
class User(BaseModel):
    user_id: str
    email: str
    name: str
    picture: Optional[str] = None
    role: Optional[str] = "parent"  # parent | teacher | student
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class Material(BaseModel):
    model_config = ConfigDict(extra="ignore")
    material_id: str = Field(default_factory=lambda: f"mat_{uuid.uuid4().hex[:10]}")
    owner: str = "guest"  # user_id or 'guest'
    title: str = "Untitled Material"
    language: str = "auto"  # zh | en | auto
    subject: str = "general"  # english | chinese | general
    text: str
    raw_ocr_text: Optional[str] = None
    ocr_confidence: Optional[str] = None  # high | low
    key_concepts: Optional[List[dict]] = None
    language_split: Optional[dict] = None
    analyzed_at: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class MaterialCreate(BaseModel):
    title: Optional[str] = "Untitled Material"
    text: str
    language: Optional[str] = "auto"
    subject: Optional[str] = "general"
    raw_ocr_text: Optional[str] = None
    ocr_confidence: Optional[str] = None


class MaterialCleanRequest(BaseModel):
    raw_text: str
    subject: Optional[str] = "general"


class MaterialCleanResponse(BaseModel):
    cleaned_text: str
    confidence: str  # high | low
    changes_summary: Optional[str] = None


class MaterialSegmentRequest(BaseModel):
    text: str
    structure: str = "word"  # word | sentence | paragraph


class MaterialSegmentResponse(BaseModel):
    formatted_text: str
    segments: List[str]
    count: int


class AnalyzeKeyConceptsRequest(BaseModel):
    text: str
    track_type: Optional[str] = "quiz"  # hint for presentation suggestions


class GenerateGameRequest(BaseModel):
    material_id: Optional[str] = None
    text: str
    mode: str  # reading_dictation | recital_dictation | quiz | exam
    subject: str = "general"  # english | chinese | general
    difficulty: int = 1  # 1-5
    exam_date: Optional[str] = None


class GameSession(BaseModel):
    session_id: str = Field(default_factory=lambda: f"gs_{uuid.uuid4().hex[:10]}")
    owner: str = "guest"
    material_id: Optional[str] = None
    track_id: Optional[str] = None
    mode: str
    score: int = 0
    max_combo: int = 0
    correct: int = 0
    wrong: int = 0
    unit_ids: List[str] = Field(default_factory=list)
    boss_defeated: bool = False
    replay_count: int = 0
    peek_count: int = 0
    completed_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class GameSessionCreate(BaseModel):
    material_id: Optional[str] = None
    track_id: Optional[str] = None
    mode: str
    score: int = 0
    max_combo: int = 0
    correct: int = 0
    wrong: int = 0
    unit_ids: List[str] = Field(default_factory=list)
    boss_defeated: bool = False
    journey_step: Optional[int] = None
    game_breakdown: Optional[dict] = None  # {G1: {correct, wrong}, ...}
    replay_count: int = 0
    peek_count: int = 0


class BehavioralSignalCreate(BaseModel):
    unit_id: Optional[str] = None
    track_id: Optional[str] = None
    game_type: Optional[str] = None
    reaction_time_ms: Optional[int] = None
    processing_time_ms: Optional[int] = None
    replay_count: int = 0
    hint_usage: int = 0
    correct: Optional[bool] = None


class StudentPreferencesUpdate(BaseModel):
    listen_mode_default: Optional[str] = None  # word | sentence


# ============== KNOWLEDGE UNIT / LEARNING TRACK (Memory Strength Engine) ==============
class KnowledgeUnit(BaseModel):
    model_config = ConfigDict(extra="ignore")
    unit_id: str = Field(default_factory=lambda: f"ku_{uuid.uuid4().hex[:10]}")
    owner: str
    track_id: str
    material_id: Optional[str] = None
    track_type: str
    unit_type: str = "word"
    term: str
    meaning: Optional[str] = None
    context: Optional[str] = None
    ease_factor: float = 2.5
    interval_days: float = 0
    repetitions: int = 0
    review_count: int = 0
    last_quality: Optional[int] = None
    last_reviewed_at: Optional[datetime] = None
    next_due_at: Optional[datetime] = None
    long_term_weak: bool = False
    language: Optional[str] = None  # zh | en — drives G4/G6/G20 variants per v3 §7.4
    key_concept_id: Optional[str] = None
    exact_term: Optional[str] = None
    simplified_explanation: Optional[str] = None
    presentation: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class LearningTrack(BaseModel):
    model_config = ConfigDict(extra="ignore")
    track_id: str = Field(default_factory=lambda: f"trk_{uuid.uuid4().hex[:10]}")
    owner: str
    student_id: Optional[str] = None  # kid who practises (guest_xxx); defaults to owner
    assigned_by: Optional[str] = None  # parent user_id when parent assigns homework
    material_id: Optional[str] = None
    track_type: str
    scope_description: str = ""
    due_date: Optional[datetime] = None
    is_cram: bool = False
    lead_time_days: Optional[int] = None
    status: str = "active"
    current_step: int = 1
    step_completed_at: dict = Field(default_factory=dict)
    consecutive_fail_count: int = 0
    source_track_id: Optional[str] = None
    current_bundle_index: int = 0
    bundle_count: int = 1
    step_lock_hours: dict = Field(default_factory=dict)
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class TrackCreate(BaseModel):
    material_id: Optional[str] = None
    text: Optional[str] = None
    title: Optional[str] = "Untitled Material"
    track_type: str
    due_date: Optional[str] = None
    student_id: Optional[str] = None  # parent assigns to linked kid (guest_xxx)
    is_template: bool = False  # teacher bulk-assign template (no single student)


class TrackAssignRequest(BaseModel):
    student_ids: Optional[List[str]] = None
    room_code: Optional[str] = None


class ProgressSnapshotSave(BaseModel):
    track_id: Optional[str] = None
    step_number: int = 1
    unit_ids: List[str] = Field(default_factory=list)
    game: dict
    progress: dict = Field(default_factory=dict)


class ReviewCreate(BaseModel):
    unit_id: str
    correct: bool
    response_time_ms: Optional[int] = None
    quality_override: Optional[int] = None


class FollowUpTask(BaseModel):
    model_config = ConfigDict(extra="ignore")
    follow_up_id: str = Field(default_factory=lambda: f"fu_{uuid.uuid4().hex[:10]}")
    owner: str
    track_id: str
    tier: str
    unit_ids: List[str]
    reviewed_unit_ids: List[str] = Field(default_factory=list)
    scheduled_date: Optional[datetime] = None
    status: str = "pending"  # pending | done | escalated_to_weak_topic
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


# ============== CONSENT / PRIVACY (ROADMAP.md Phase 2 — COPPA / GDPR-K / PDPO) ==============
# This is not a nice-to-have: it's the prerequisite for retaining ANY identifiable
# data about a child beyond a single guest session. Nothing else in this file reads
# these collections yet — that's intentional wiring for the next step, not a bug —
# but the data model needs to exist now so it's not bolted on as an afterthought
# once real users are already in the database.

PRIVACY_POLICY_VERSION = "2026-07-v1"  # bump this every time the policy text changes


class FamilyLink(BaseModel):
    model_config = ConfigDict(extra="ignore")
    link_id: str = Field(default_factory=lambda: f"fl_{uuid.uuid4().hex[:10]}")
    parent_user_id: str = ""
    kid_owner_id: str  # guest_{kid_device_id}
    kid_device_id: Optional[str] = None
    permission_level: str = "manage_tracks"  # "view_full_analysis" | "manage_tracks"
    status: str = "active"  # "pending" | "active"
    parent_email: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class FamilyLinkCreate(BaseModel):
    kid_owner_id: str
    permission_level: Optional[str] = "manage_tracks"


class FamilyLinkByCodeCreate(BaseModel):
    family_code: str
    permission_level: Optional[str] = "manage_tracks"


class FamilyLinkInviteCreate(BaseModel):
    parent_email: str
    kid_device_id: Optional[str] = None


class ConsentRecord(BaseModel):
    model_config = ConfigDict(extra="ignore")
    consent_id: str = Field(default_factory=lambda: f"cn_{uuid.uuid4().hex[:10]}")
    parent_user_id: str = ""
    kid_owner_id: str
    parent_email: Optional[str] = None
    link_id: Optional[str] = None
    # Kept as SEPARATE consent types on purpose — COPPA/PDPO best practice is not to
    # bundle "let my kid use the app" with "share aggregated data with a research
    # partner" under one checkbox. A parent can grant one without the other.
    consent_type: str  # "data_collection" | "research_sharing" | "marketing"
    granted: bool = False
    method: str = "email_link"  # "email_link" | "in_app_signature"
    policy_version: str = PRIVACY_POLICY_VERSION
    consent_token: Optional[str] = None  # set while pending, cleared once confirmed
    granted_at: Optional[datetime] = None
    revoked_at: Optional[datetime] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class ConsentRequest(BaseModel):
    kid_owner_id: str
    consent_type: str
    parent_email: str


async def _send_consent_email(parent_email: str, confirm_url: str, consent_type: str):
    await send_consent_email(parent_email, confirm_url, consent_type)


# ============== AUTH ==============
async def get_current_user(
    request: Request,
    session_token: Optional[str] = Cookie(default=None),
    authorization: Optional[str] = Header(default=None),
) -> Optional[User]:
    token = session_token
    if not token and authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ", 1)[1]
    if not token:
        return None
    sess = await db.user_sessions.find_one({"session_token": token}, {"_id": 0})
    if not sess:
        return None
    expires_at = sess.get("expires_at")
    if isinstance(expires_at, str):
        expires_at = datetime.fromisoformat(expires_at)
    if expires_at and expires_at.tzinfo is None:
        expires_at = expires_at.replace(tzinfo=timezone.utc)
    if expires_at and expires_at < datetime.now(timezone.utc):
        return None
    user_doc = await db.users.find_one({"user_id": sess["user_id"]}, {"_id": 0})
    if not user_doc:
        return None
    return User(**user_doc)


async def get_owner_id(
    user: Optional[User] = Depends(get_current_user),
    x_guest_id: Optional[str] = Header(default=None, alias="X-Guest-Id"),
    x_kid_mode: Optional[str] = Header(default=None, alias="X-Kid-Mode"),
) -> str:
    # Resolve a stable owner id: logged-in user first, else a per-device guest id
    # sent by the frontend as X-Guest-Id (generated once, stored in localStorage).
    # When X-Kid-Mode is set (kid hub / upload / battle), always use the device
    # guest id even if a parent/teacher session cookie is still present — otherwise
    # "Start Adventure" after visiting the guardian portal writes to the wrong owner.
    kid_mode = str(x_kid_mode or "").lower() in ("1", "true", "yes")
    if kid_mode and x_guest_id:
        return f"guest_{x_guest_id.removeprefix('guest_')}"
    if user:
        return user.user_id
    if x_guest_id:
        return f"guest_{x_guest_id.removeprefix('guest_')}"
    return "guest_anonymous"


def kid_owner_from_device(kid_device_id: str) -> str:
    raw = kid_device_id.removeprefix("guest_")
    return f"guest_{raw}"


def _student_tracks_filter(student_owner: str) -> dict:
    """Tracks visible in a kid's daily queue (own + parent-assigned)."""
    return {
        "status": "active",
        "$or": [
            {"student_id": student_owner},
            {"owner": student_owner, "$or": [{"student_id": None}, {"student_id": {"$exists": False}}]},
        ],
    }


async def _gen_family_code() -> str:
    for _ in range(30):
        code = f"{random.randint(0, 999999):06d}"
        if not await db.family_codes.find_one({"code": code}, {"_id": 0}):
            return code
    raise HTTPException(status_code=500, detail="Could not generate family code")


async def _ensure_family_code(kid_device_id: str) -> str:
    raw = kid_device_id.removeprefix("guest_")
    existing = await db.family_codes.find_one({"kid_device_id": raw}, {"_id": 0})
    if existing:
        return existing["code"]
    code = await _gen_family_code()
    await db.family_codes.insert_one({
        "code": code,
        "kid_device_id": raw,
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    return code


async def _resolve_family_code(code: str) -> dict:
    doc = await db.family_codes.find_one({"code": code.strip()}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Family code not found")
    return doc


async def _active_family_link_count(kid_owner_id: str) -> int:
    return await db.family_links.count_documents({"kid_owner_id": kid_owner_id, "status": "active"})


async def _parent_linked_kid_ids(parent_user_id: str) -> list:
    links = await db.family_links.find(
        {"parent_user_id": parent_user_id, "status": "active"}, {"_id": 0, "kid_owner_id": 1}
    ).to_list(50)
    return [l["kid_owner_id"] for l in links if l.get("kid_owner_id")]


async def _resolve_child_data_owner(
    user: Optional[User],
    owner: str,
    kid_owner_id: Optional[str] = None,
) -> Optional[str]:
    """Parent portal: only return data for a linked child. None => empty payload."""
    if not user or owner != user.user_id:
        return owner
    kid_ids = await _parent_linked_kid_ids(user.user_id)
    if not kid_ids:
        return None
    if kid_owner_id:
        if kid_owner_id not in kid_ids:
            raise HTTPException(status_code=403, detail="Not linked to this child")
        return kid_owner_id
    if len(kid_ids) == 1:
        return kid_ids[0]
    raise HTTPException(status_code=400, detail="kid_owner_id required when multiple children linked")


async def _require_parent_child_link(user: User, kid_owner_id: str) -> None:
    link = await db.family_links.find_one(
        {"parent_user_id": user.user_id, "kid_owner_id": kid_owner_id, "status": "active"},
        {"_id": 0},
    )
    if not link:
        raise HTTPException(status_code=403, detail="Not linked to this child")


async def _resolve_track_student(
    owner: str,
    user: Optional[User],
    payload: "TrackCreate",
) -> tuple[str, Optional[str]]:
    """Return (student_id, assigned_by) for a new track."""
    if payload.is_template:
        return None, user.user_id if user else None

    if payload.track_type == "self_practice":
        if user:
            return owner, None
        return owner, None

    if user:
        target = payload.student_id
        if target:
            link = await db.family_links.find_one(
                {"parent_user_id": user.user_id, "kid_owner_id": target, "status": "active"},
                {"_id": 0},
            )
            if not link:
                raise HTTPException(status_code=403, detail="Not linked to this child")
            return target, user.user_id
        links = await db.family_links.find(
            {"parent_user_id": user.user_id, "status": "active"}, {"_id": 0}
        ).to_list(20)
        if links:
            if len(links) == 1:
                return links[0]["kid_owner_id"], user.user_id
            raise HTTPException(
                status_code=400,
                detail="Multiple children linked — pass student_id when assigning",
            )
        return owner, None

    return owner, None


def _track_access_filter(owner: str) -> dict:
    return {
        "$or": [
            {"owner": owner},
            {"student_id": owner},
        ],
    }


class AuthSessionCreate(BaseModel):
    id_token: str
    role: str = "parent"


class DevLoginCreate(BaseModel):
    email: str
    name: Optional[str] = None
    role: str = "parent"


async def _create_user_session(response: Response, email: str, name: str, picture: Optional[str], role: str):
    existing = await db.users.find_one({"email": email}, {"_id": 0})
    if existing:
        user_id = existing["user_id"]
        await db.users.update_one(
            {"user_id": user_id},
            {"$set": {"name": name, "picture": picture, "role": role}},
        )
    else:
        user_id = f"user_{uuid.uuid4().hex[:12]}"
        await db.users.insert_one({
            "user_id": user_id,
            "email": email,
            "name": name,
            "picture": picture,
            "role": role,
            "created_at": datetime.now(timezone.utc).isoformat(),
        })

    session_token = f"sess_{uuid.uuid4().hex}"
    expires_at = datetime.now(timezone.utc) + timedelta(days=7)
    await db.user_sessions.insert_one({
        "user_id": user_id,
        "session_token": session_token,
        "expires_at": expires_at.isoformat(),
        "created_at": datetime.now(timezone.utc).isoformat(),
    })

    is_local = os.environ.get("APP_ENV", "development") == "development"
    response.set_cookie(
        key="session_token",
        value=session_token,
        max_age=7 * 24 * 60 * 60,
        httponly=True,
        secure=not is_local,
        samesite="lax" if is_local else "none",
        path="/",
    )

    return {
        "user_id": user_id,
        "email": email,
        "name": name,
        "picture": picture,
        "role": role,
    }


@api_router.post("/auth/session")
async def auth_session(payload: AuthSessionCreate, response: Response):
    profile = await verify_firebase_id_token(payload.id_token)
    if not profile or not profile.get("email"):
        raise HTTPException(status_code=401, detail="Invalid Firebase id_token")
    return await _create_user_session(
        response,
        profile["email"],
        profile["name"],
        profile.get("picture"),
        payload.role,
    )


@api_router.post("/auth/dev-login")
async def dev_login(payload: DevLoginCreate, response: Response):
    if os.environ.get("DEV_AUTH_ENABLED", "").lower() != "true":
        raise HTTPException(status_code=403, detail="Dev auth disabled")
    name = payload.name or payload.email.split("@")[0]
    return await _create_user_session(response, payload.email, name, None, payload.role)


@api_router.get("/auth/me")
async def auth_me(user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Not authenticated")
    return user


@api_router.post("/auth/logout")
async def auth_logout(
    response: Response,
    session_token: Optional[str] = Cookie(default=None),
):
    if session_token:
        await db.user_sessions.delete_one({"session_token": session_token})
    response.delete_cookie("session_token", path="/")
    return {"ok": True}


# ============== MATERIALS ==============

_CJK_CHAR = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf]")
_CJK_SPACE_CJK = re.compile(r"([\u4e00-\u9fff\u3400-\u4dbf])\s+([\u4e00-\u9fff\u3400-\u4dbf])")
_EN_WORD = re.compile(r"[a-zA-Z]{3,}")


def _clean_ocr_fallback(raw_text: str) -> tuple[str, str]:
    """Deterministic OCR cleanup when LLM is unavailable."""
    text = raw_text or ""
    # Remove spaces between consecutive CJK characters (keep CJK↔Latin gaps).
    prev = None
    while prev != text:
        prev = text
        text = _CJK_SPACE_CJK.sub(r"\1\2", text)
    # Collapse 3+ blank lines into one paragraph break.
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            lines.append("")
            continue
        has_cjk = bool(_CJK_CHAR.search(stripped))
        has_en = bool(_EN_WORD.search(stripped))
        if len(stripped) < 20 and not has_cjk and not has_en:
            continue
        lines.append(stripped)
    cleaned = "\n".join(lines).strip()
    summary = "Fallback cleanup: removed CJK spacing noise and junk lines"
    return cleaned or raw_text.strip(), summary


_OCR_CLEAN_SYSTEM = """You are an OCR post-processing assistant. Your ONLY job is to clean noise from OCR output — NOT rewrite, NOT fill in answers, NOT invent content.

MUST DO:
- Remove extra spaces between Chinese characters ("小 五 中 文" → "小五中文")
- Fix obvious line breaks that break sentences into fragments
- Remove meaningless garbage fragments (handwritten signatures, page numbers, random alphanumeric noise like "29Junim 991 XS") — note in changes_summary when removed

NEVER DO:
- Do NOT fill in blank answers or invent missing content
- Do NOT change numbers or factual content (if "$2.82" might be wrong, note suspicion in changes_summary but do NOT substitute another number — use [?] if needed)
- Do NOT add information not in the original

If a section is too garbled to understand, keep it as-is and say in changes_summary that OCR quality is too low and the user should re-photograph.

Respond with ONLY valid JSON (no markdown fences):
{"cleaned_text": "...", "changes_summary": "..."}"""


@api_router.post("/materials/clean", response_model=MaterialCleanResponse)
async def clean_material_text(payload: MaterialCleanRequest):
    raw = (payload.raw_text or "").strip()
    if not raw:
        raise HTTPException(status_code=400, detail="raw_text is required")
    if len(raw) < 4:
        return MaterialCleanResponse(cleaned_text=raw, confidence="low", changes_summary="Text too short")

    if GEMINI_ENABLED:
        try:
            prompt = f"Subject context: {payload.subject or 'general'}\n\nOCR raw text:\n{raw[:12000]}"
            out = await _gemini_generate(_OCR_CLEAN_SYSTEM, prompt, timeout=9.0)
            m = re.search(r"\{[\s\S]*\}", out)
            if m:
                data = json.loads(m.group())
                cleaned = (data.get("cleaned_text") or "").strip()
                if cleaned:
                    return MaterialCleanResponse(
                        cleaned_text=cleaned,
                        confidence="high",
                        changes_summary=data.get("changes_summary"),
                    )
        except Exception as exc:
            logger.warning("materials/clean LLM failed: %s", exc)

    cleaned, summary = _clean_ocr_fallback(raw)
    return MaterialCleanResponse(cleaned_text=cleaned, confidence="high", changes_summary=summary)


_SEGMENT_SYSTEM = """You split Chinese/English study material into learning items.
structure=word: vocabulary, idioms, terms — ONE item per line, exact characters only, no invented content.
  For run-on Chinese with no spaces (e.g. 春眠不覺曉處處聞啼鳥), split into 2–4 character words/phrases as students would learn them.
structure=sentence: one sentence per line (split on 。！？.!?).
structure=paragraph: one paragraph per line.

Return ONLY JSON: {"segments": ["item1", "item2", ...]}"""


def _segment_text_deterministic(text: str, structure: str) -> List[str]:
    text = (text or "").strip()
    if not text:
        return []
    if structure == "sentence":
        parts = re.split(r"(?<=[。！？.!?])\s*|\n+", text)
        return [p.strip() for p in parts if len(p.strip()) >= 2]
    if structure == "paragraph":
        parts = re.split(r"\n\s*\n+", text)
        if len(parts) > 1:
            return [p.strip() for p in parts if len(p.strip()) >= 5]
        # Long single block — split at major punctuation every ~80 chars
        chunks, buf = [], ""
        for ch in text:
            buf += ch
            if len(buf) >= 80 and ch in "。！？":
                chunks.append(buf.strip())
                buf = ""
        if buf.strip():
            chunks.append(buf.strip())
        return chunks if chunks else [text]
    # word / terms — split on punctuation and whitespace first
    parts = re.split(r"[，、；。\s]+", text)
    segments = [p.strip() for p in parts if len(p.strip()) >= 1]
    if len(segments) <= 1 and len(text) > 6 and re.search(r"[\u4e00-\u9fff]", text):
        return []  # signal: run-on CJK blob needs LLM
    return segments


def _segment_cjk_vocabulary_greedy(text: str) -> List[str]:
    """Split run-on Chinese vocabulary into 1–4 character terms (no LLM)."""
    text = re.sub(r"\s+", "", text or "")
    if not text:
        return []
    for sep in ("、", "，", "；", "。", "·", "|"):
        if sep in text:
            parts = [p.strip() for p in text.split(sep) if p.strip()]
            if len(parts) > 1:
                out: List[str] = []
                for p in parts:
                    sub = _segment_cjk_vocabulary_greedy(p)
                    out.extend(sub if sub else [p])
                flat = [w for w in out if w.strip()]
                if len(flat) > 1:
                    return flat
    words: List[str] = []
    i = 0
    while i < len(text):
        ch = text[i]
        if not ("\u4e00" <= ch <= "\u9fff"):
            i += 1
            continue
        placed = False
        for n in (4, 3, 2, 1):
            chunk = text[i : i + n]
            if len(chunk) == n and all("\u4e00" <= c <= "\u9fff" for c in chunk):
                words.append(chunk)
                i += n
                placed = True
                break
        if not placed:
            words.append(ch)
            i += 1
    return [w for w in words if w.strip()]


@api_router.post("/materials/segment", response_model=MaterialSegmentResponse)
async def segment_material_text(payload: MaterialSegmentRequest):
    raw = (payload.text or "").strip()
    structure = (payload.structure or "word").lower()
    if structure not in ("word", "sentence", "paragraph"):
        structure = "word"
    if len(raw) < 2:
        raise HTTPException(status_code=400, detail="Need text to segment")

    segments = _segment_text_deterministic(raw, structure)

    if (not segments or (structure == "word" and len(segments) == 1 and len(segments[0]) > 6)) and GEMINI_ENABLED:
        try:
            prompt = f"structure: {structure}\n\nMaterial:\n{raw[:8000]}"
            out = await _gemini_generate(_SEGMENT_SYSTEM, prompt, timeout=12.0)
            m = re.search(r"\{[\s\S]*\}", out)
            if m:
                data = json.loads(m.group())
                llm_segs = [s.strip() for s in (data.get("segments") or []) if isinstance(s, str) and s.strip()]
                if llm_segs:
                    segments = llm_segs
        except Exception as exc:
            logger.warning("materials/segment LLM failed: %s", exc)

    if structure == "word" and (not segments or (len(segments) == 1 and len(segments[0]) > 8)):
        if re.search(r"[\u4e00-\u9fff]", raw):
            greedy = _segment_cjk_vocabulary_greedy(raw)
            if len(greedy) > 1:
                segments = greedy

    if not segments:
        segments = _segment_text_deterministic(raw, "sentence") or [raw]

    formatted = "\n".join(segments)
    return MaterialSegmentResponse(formatted_text=formatted, segments=segments, count=len(segments))


@api_router.post("/materials")
async def create_material(
    payload: MaterialCreate,
    owner: str = Depends(get_owner_id),
):
    await assert_can_collect_learning_data(db, owner)
    mat = Material(
        owner=owner,
        title=payload.title or "Untitled Material",
        text=payload.text,
        language=payload.language or "auto",
        subject=payload.subject or "general",
        raw_ocr_text=payload.raw_ocr_text,
        ocr_confidence=payload.ocr_confidence,
    )
    doc = mat.model_dump()
    doc["created_at"] = doc["created_at"].isoformat()
    await db.materials.insert_one(doc)
    return mat


@api_router.get("/materials")
async def list_materials(
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
    kid_owner_id: Optional[str] = None,
):
    data_owner = await _resolve_child_data_owner(user, owner, kid_owner_id)
    if data_owner is None:
        return []
    docs = await db.materials.find({"owner": data_owner}, {"_id": 0}).sort("created_at", -1).to_list(50)
    return docs


@api_router.post("/materials/analyze-key-concepts")
async def analyze_key_concepts(payload: AnalyzeKeyConceptsRequest, owner: str = Depends(get_owner_id)):
    """§7.1 Key concept detection after upload — preview before track creation."""
    text = (payload.text or "").strip()
    if len(text) < 8:
        raise HTTPException(status_code=400, detail="Need more text to analyze")
    result = await extract_key_concepts(
        text,
        payload.track_type or "quiz",
        gemini_generate=_gemini_generate,
        gemini_enabled=GEMINI_ENABLED,
    )
    return result


class OcrRequest(BaseModel):
    image_base64: str
    lang_hint: Optional[str] = "auto"  # auto | eng | zh | zh_trad | zh_simp
    engine: Optional[str] = "auto"  # auto | gemini | vision


class DeviceProfileUpdate(BaseModel):
    subscription_tier: Optional[str] = None  # free | premium


class ChildSubscriptionUpdate(BaseModel):
    kid_owner_id: str
    subscription_tier: str  # free | premium


class DisplayNameUpdate(BaseModel):
    display_name: str


class FriendAddRequest(BaseModel):
    public_user_id: str


_PROFANITY_WORDS = {
    "fuck", "shit", "damn", "bitch", "asshole", "bastard", "cunt", "dick", "pussy",
    "操", "傻逼", "傻B", "妈的", "他媽", "他妈", "屌", "撚", "撚樣", "仆街", "死全家",
}


def _contains_profanity(name: str) -> bool:
    lower = (name or "").lower()
    for w in _PROFANITY_WORDS:
        if w.lower() in lower:
            return True
    return False


def _gen_public_user_id() -> str:
    import string
    chars = string.ascii_uppercase + string.digits
    return "".join(random.choices(chars, k=8))


async def _ensure_public_profile(owner: str) -> dict:
    profile = await get_device_profile(db, owner)
    updates = {}
    if not profile.get("public_user_id"):
        for _ in range(20):
            pid = _gen_public_user_id()
            clash = await db.device_profiles.find_one({"public_user_id": pid}, {"_id": 1})
            if not clash:
                updates["public_user_id"] = pid
                break
    if not profile.get("display_name"):
        updates["display_name"] = "Agent"
    if updates:
        updates["updated_at"] = datetime.now(timezone.utc).isoformat()
        await db.device_profiles.update_one({"owner": owner}, {"$set": updates}, upsert=True)
        profile = await get_device_profile(db, owner)
    return profile


@api_router.get("/ocr/quota")
async def ocr_quota(owner: str = Depends(get_owner_id)):
    return await safe_get_quota_status(db, owner)


@api_router.get("/ocr/capabilities")
async def ocr_capabilities():
    return {
        "gemini": GEMINI_ENABLED,
        "vision": VISION_ENABLED,
        "tesseract": True,
        "mlkit": False,
        "languages": ["auto", "eng", "zh", "zh_trad", "zh_simp"],
        "daily_limit_free": int(os.environ.get("AI_SCANS_PER_DAY_FREE", "3")),
    }


@api_router.post("/ocr")
async def ocr_image(
    payload: OcrRequest,
    owner: str = Depends(get_owner_id),
    x_user_gemini_key: Optional[str] = Header(default=None, alias="X-User-Gemini-Key"),
):
    """Hybrid OCR: Vision → Gemini (server or BYOK). Counts against daily quota unless BYOK."""
    try:
        image_bytes, mime = decode_image_base64(payload.image_base64)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid image data")
    if len(image_bytes) > 8 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="Image too large (max 8MB)")
    byok = (x_user_gemini_key or "").strip()
    if byok.startswith("PASTE"):
        byok = ""
    try:
        result = await run_hybrid_ocr(
            db, owner, image_bytes, mime, payload.lang_hint or "auto",
            engine=payload.engine or "auto",
            skip_quota=bool(byok),
            byok_gemini_key=byok or None,
            gemini_enabled=GEMINI_ENABLED,
            client_factory=_get_genai_client,
            genai_types=genai_types,
        )
    except QuotaExceededError as e:
        raise HTTPException(
            status_code=402,
            detail={"message": "AI scan quota exceeded — use basic OCR or paste text", "quota": e.quota},
        )
    except RateLimitError:
        raise HTTPException(
            status_code=429,
            detail="AI servers busy — try again later or use basic OCR / paste text",
        )
    except asyncio.TimeoutError:
        raise HTTPException(status_code=504, detail="OCR timed out")
    except RuntimeError as e:
        if "not set" in str(e).lower() or "not configured" in str(e).lower():
            raise HTTPException(status_code=503, detail=str(e))
        raise HTTPException(status_code=502, detail=str(e))
    if not result.get("text"):
        raise HTTPException(status_code=422, detail="No readable text found")
    return result


# ============== DOCUMENT UPLOAD (PDF / Word / PowerPoint — no OCR needed, real embedded text) ==============
DOC_MAX_SIZE_MB = int(os.environ.get("DOC_MAX_SIZE_MB", "20"))
# Deliberately capped well below "300 pages" (Section: user worried about kids
# uploading a whole textbook) — this isn't a hard block, we just read the first
# N pages and tell the UI it was truncated, so the daily learning queue never
# balloons into something no child could realistically get through anyway.
DOC_MAX_PAGES_FREE = int(os.environ.get("DOC_MAX_PAGES_FREE", "15"))
DOC_MAX_PAGES_PREMIUM = int(os.environ.get("DOC_MAX_PAGES_PREMIUM", "40"))
DOC_MAX_FILES_FREE = int(os.environ.get("DOC_MAX_FILES_FREE", "3"))
DOC_MAX_FILES_PREMIUM = int(os.environ.get("DOC_MAX_FILES_PREMIUM", "8"))


def _extract_pdf_text(raw: bytes, max_pages: int):
    import io
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    total_pages = len(reader.pages)
    pages_used = min(total_pages, max_pages)
    parts = []
    for i in range(pages_used):
        try:
            parts.append(reader.pages[i].extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(p.strip() for p in parts if p.strip()), total_pages, pages_used


def _extract_docx_text(raw: bytes, max_pages: int):
    import io
    from docx import Document

    doc = Document(io.BytesIO(raw))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    # .docx has no real "page" concept at the XML level — approximate a page as
    # ~40 paragraphs so long Word docs respect the same sane cap as PDFs/slides.
    approx_page_size = 40
    total_pages = max(1, -(-len(paras) // approx_page_size))
    pages_used = min(total_pages, max_pages)
    kept = paras[: pages_used * approx_page_size]
    return "\n".join(kept), total_pages, pages_used


def _extract_pptx_text(raw: bytes, max_pages: int):
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    slides = list(prs.slides)
    total_pages = len(slides)
    pages_used = min(total_pages, max_pages)
    parts = []
    for slide in slides[:pages_used]:
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text
                if t.strip():
                    texts.append(t.strip())
        if texts:
            parts.append("\n".join(texts))
    return "\n\n".join(parts), total_pages, pages_used


@api_router.get("/documents/limits")
async def document_limits(owner: str = Depends(get_owner_id)):
    profile = await get_device_profile(db, owner)
    is_premium = profile.get("subscription_tier") == "premium"
    return {
        "subscription_tier": profile.get("subscription_tier", "free"),
        "max_files": DOC_MAX_FILES_PREMIUM if is_premium else DOC_MAX_FILES_FREE,
        "max_pages_per_file": DOC_MAX_PAGES_PREMIUM if is_premium else DOC_MAX_PAGES_FREE,
        "max_file_size_mb": DOC_MAX_SIZE_MB,
    }


@api_router.post("/documents/extract")
async def extract_document(
    file: UploadFile = File(...),
    owner: str = Depends(get_owner_id),
):
    """Extract text from an uploaded PDF/Word/PowerPoint file — these have real
    embedded text so no OCR/AI-scan-quota is needed at all. Pages are capped
    per subscription tier (see DOC_MAX_PAGES_*); we still return whatever we
    could read plus a `truncated` flag rather than hard-rejecting long files."""
    filename = (file.filename or "").lower()
    raw = await file.read()
    size_mb = len(raw) / (1024 * 1024)
    if size_mb > DOC_MAX_SIZE_MB:
        raise HTTPException(status_code=400, detail=f"File too large (max {DOC_MAX_SIZE_MB}MB)")

    profile = await get_device_profile(db, owner)
    is_premium = profile.get("subscription_tier") == "premium"
    max_pages = DOC_MAX_PAGES_PREMIUM if is_premium else DOC_MAX_PAGES_FREE

    try:
        if filename.endswith(".pdf"):
            text, total_pages, pages_used = _extract_pdf_text(raw, max_pages)
            kind = "pdf"
        elif filename.endswith(".docx"):
            text, total_pages, pages_used = _extract_docx_text(raw, max_pages)
            kind = "docx"
        elif filename.endswith(".pptx"):
            text, total_pages, pages_used = _extract_pptx_text(raw, max_pages)
            kind = "pptx"
        elif filename.endswith(".doc") or filename.endswith(".ppt"):
            raise HTTPException(
                status_code=415,
                detail="Old .doc/.ppt format isn't supported — please save/export as .docx or .pptx",
            )
        else:
            raise HTTPException(
                status_code=415,
                detail="Unsupported file type — use PDF, Word (.docx) or PowerPoint (.pptx)",
            )
    except HTTPException:
        raise
    except ImportError:
        raise HTTPException(
            status_code=501,
            detail="Document parsing isn't set up on the server yet (pypdf/python-docx/python-pptx missing)",
        )
    except Exception as e:
        logger.warning("document extraction failed for %s: %s", filename, e)
        raise HTTPException(
            status_code=422,
            detail="Couldn't read this file — it may be corrupted, scanned-image-only, or password protected",
        )

    if not text.strip():
        raise HTTPException(
            status_code=422,
            detail="No extractable text found — if this is a scanned/image-only document, take photos of the pages instead so AI OCR can read them",
        )

    return {
        "text": text.strip(),
        "source": kind,
        "filename": file.filename,
        "total_pages": total_pages,
        "pages_used": pages_used,
        "truncated": pages_used < total_pages,
        "max_pages": max_pages,
        "subscription_tier": profile.get("subscription_tier", "free"),
    }


@api_router.patch("/device-profile")
async def update_device_profile(
    payload: DeviceProfileUpdate,
    owner: str = Depends(get_owner_id),
):
    """Set subscription_tier (premium unlocks unlimited AI scans — IAP wiring later)."""
    updates = {"updated_at": datetime.now(timezone.utc).isoformat()}
    if payload.subscription_tier in ("free", "premium"):
        updates["subscription_tier"] = payload.subscription_tier
    await db.device_profiles.update_one({"owner": owner}, {"$set": updates}, upsert=True)
    return await get_device_profile(db, owner)


@api_router.get("/device-profile")
async def read_device_profile(owner: str = Depends(get_owner_id)):
    return await get_device_profile(db, owner)


@api_router.get("/profile/me")
async def get_profile_me(owner: str = Depends(get_owner_id)):
    profile = await _ensure_public_profile(owner)
    friends = await db.friends.find({"owner": owner}, {"_id": 0, "friend_owner": 1, "friend_public_id": 1, "friend_display_name": 1}).to_list(50)
    return {
        "owner": owner,
        "display_name": profile.get("display_name", "Agent"),
        "public_user_id": profile.get("public_user_id"),
        "friend_count": len(friends),
    }


@api_router.patch("/profile/display-name")
async def update_display_name(payload: DisplayNameUpdate, owner: str = Depends(get_owner_id)):
    name = (payload.display_name or "").strip()
    if len(name) < 1 or len(name) > 24:
        raise HTTPException(status_code=400, detail="Display name must be 1–24 characters")
    if _contains_profanity(name):
        raise HTTPException(status_code=400, detail="Display name contains inappropriate words")
    await _ensure_public_profile(owner)
    await db.device_profiles.update_one(
        {"owner": owner},
        {"$set": {"display_name": name, "updated_at": datetime.now(timezone.utc).isoformat()}},
        upsert=True,
    )
    return {"display_name": name}


@api_router.get("/users/search")
async def search_user_by_id(public_user_id: str, owner: str = Depends(get_owner_id)):
    pid = (public_user_id or "").strip().upper()
    if len(pid) < 4:
        raise HTTPException(status_code=400, detail="Enter a valid user ID")
    doc = await db.device_profiles.find_one({"public_user_id": pid}, {"_id": 0, "owner": 1, "display_name": 1, "public_user_id": 1})
    if not doc:
        raise HTTPException(status_code=404, detail="User not found")
    if doc["owner"] == owner:
        raise HTTPException(status_code=400, detail="That's your own ID")
    return {
        "public_user_id": doc.get("public_user_id"),
        "display_name": doc.get("display_name", "Agent"),
    }


@api_router.post("/friends")
async def add_friend(payload: FriendAddRequest, owner: str = Depends(get_owner_id)):
    pid = (payload.public_user_id or "").strip().upper()
    target = await db.device_profiles.find_one({"public_user_id": pid}, {"_id": 0})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    friend_owner = target["owner"]
    if friend_owner == owner:
        raise HTTPException(status_code=400, detail="Cannot add yourself")
    existing = await db.friends.find_one({"owner": owner, "friend_owner": friend_owner}, {"_id": 0})
    if existing:
        return {"ok": True, "already_friends": True}
    doc = {
        "friend_id": f"fr_{uuid.uuid4().hex[:10]}",
        "owner": owner,
        "friend_owner": friend_owner,
        "friend_public_id": pid,
        "friend_display_name": target.get("display_name", "Agent"),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.friends.insert_one(doc)
    return {"ok": True, "friend": doc}


@api_router.get("/friends")
async def list_friends(owner: str = Depends(get_owner_id)):
    docs = await db.friends.find({"owner": owner}, {"_id": 0}).sort("created_at", -1).to_list(50)
    out = []
    for f in docs:
        prof = await db.device_profiles.find_one({"owner": f["friend_owner"]}, {"_id": 0, "display_name": 1, "public_user_id": 1})
        sessions = await db.game_sessions.find({"owner": f["friend_owner"]}, {"_id": 0, "score": 1}).to_list(200)
        xp = sum(s.get("score", 0) for s in sessions)
        out.append({
            "friend_id": f.get("friend_id"),
            "display_name": prof.get("display_name") if prof else f.get("friend_display_name", "Agent"),
            "public_user_id": prof.get("public_user_id") if prof else f.get("friend_public_id"),
            "total_xp": xp,
        })
    me = await _ensure_public_profile(owner)
    my_sessions = await db.game_sessions.find({"owner": owner}, {"_id": 0, "score": 1}).to_list(200)
    my_xp = sum(s.get("score", 0) for s in my_sessions)
    return {
        "me": {"display_name": me.get("display_name", "Agent"), "public_user_id": me.get("public_user_id"), "total_xp": my_xp},
        "friends": out,
    }


# ============== AVATAR / DIAMOND SKIN SHOP (Section 3: cosmetic-only, never gameplay) ==============
AVATAR_CATALOG = [
    {"skin_id": "fox", "emoji": "🦊", "name": "Fox", "cost": 0},
    {"skin_id": "owl", "emoji": "🦉", "name": "Owl", "cost": 0},
    {"skin_id": "robot", "emoji": "🤖", "name": "Robot", "cost": 0},
    {"skin_id": "dragon", "emoji": "🐉", "name": "Dragon", "cost": 20},
    {"skin_id": "unicorn", "emoji": "🦄", "name": "Unicorn", "cost": 20},
    {"skin_id": "ninja", "emoji": "🥷", "name": "Ninja", "cost": 30},
    {"skin_id": "wizard", "emoji": "🧙", "name": "Wizard", "cost": 30},
    {"skin_id": "astronaut", "emoji": "🧑\u200d🚀", "name": "Astronaut", "cost": 40},
    {"skin_id": "phoenix", "emoji": "🐦\u200d🔥", "name": "Phoenix", "cost": 60},
]
DEFAULT_UNLOCKED_SKINS = [s["skin_id"] for s in AVATAR_CATALOG if s["cost"] == 0]


class AvatarSkinRequest(BaseModel):
    skin_id: str


@api_router.get("/student-preferences")
async def get_student_preferences(owner: str = Depends(get_owner_id)):
    doc = await db.student_preferences.find_one({"student_id": owner}, {"_id": 0})
    return doc or {"student_id": owner, "listen_mode_default": "word"}


@api_router.post("/student-preferences")
async def save_student_preferences(payload: StudentPreferencesUpdate, owner: str = Depends(get_owner_id)):
    now = datetime.now(timezone.utc).isoformat()
    update = {k: v for k, v in payload.model_dump().items() if v is not None}
    update["student_id"] = owner
    update["updated_at"] = now
    await db.student_preferences.update_one({"student_id": owner}, {"$set": update}, upsert=True)
    return update


@api_router.get("/profile/avatar")
async def get_avatar_profile(owner: str = Depends(get_owner_id)):
    profile = await get_device_profile(db, owner)
    unlocked = profile.get("unlocked_skins") or DEFAULT_UNLOCKED_SKINS
    equipped = profile.get("equipped_skin") or "fox"
    badges = await get_student_badges(db, owner)
    return {
        "diamonds": profile.get("diamonds", 0),
        "coins": profile.get("coins", 0),
        "unlocked_skins": unlocked,
        "equipped_skin": equipped,
        "catalog": AVATAR_CATALOG,
        "badges": badges,
        "badge_catalog": BADGE_CATALOG,
    }


@api_router.get("/badges")
async def list_badges(owner: str = Depends(get_owner_id)):
    return {"badges": await get_student_badges(db, owner), "catalog": BADGE_CATALOG}


# ============== G17 AI SENTENCE GRADING (v3 §7.3) ==============
GRADE_SENTENCE_SYSTEM = """You grade a primary-school child's sentence for correct use of a keyword.
Return ONLY valid JSON with keys: passed (bool), score_pct (0-100 int), feedback (short encouraging hint in the same language as the student's sentence).
Rules:
- The ONLY keyword to check is the one given in "Keyword:" — never mention a different word in feedback.
- passed=true if the keyword is used correctly in context OR the answer captures the right concept direction.
- passed=false if keyword missing or used wrongly.
- feedback must reference the exact keyword given, gentle, never harsh."""


class GradeSentenceRequest(BaseModel):
    keyword: str
    sentence: str
    meaning: Optional[str] = None


def _grade_sentence_fallback(keyword: str, sentence: str, meaning: Optional[str] = None) -> dict:
    kw = (keyword or "").strip()
    text = (sentence or "").strip()
    if not text:
        return {"passed": False, "score_pct": 0, "feedback": "請寫一句完整句子。"}
    if kw and kw in text:
        return {"passed": True, "score_pct": 100, "feedback": "用字正確，做得好！"}
    if meaning and any(part in text for part in meaning.replace("，", " ").split() if len(part) >= 2):
        return {"passed": True, "score_pct": 75, "feedback": f"意思啱，記得用返「{kw}」會更好。"}
    return {
        "passed": False,
        "score_pct": 25,
        "feedback": f"句子要用到「{kw}」先算啱。" if kw else "再試一次，用返指定字詞。",
    }


@api_router.post("/grade-sentence")
async def grade_sentence(payload: GradeSentenceRequest, owner: str = Depends(get_owner_id)):
    kw = (payload.keyword or "").strip()
    sentence = (payload.sentence or "").strip()
    if not sentence:
        raise HTTPException(status_code=400, detail="sentence required")
    result = None
    if GEMINI_ENABLED and kw:
        try:
            prompt = f"Keyword: {kw}\nMeaning: {payload.meaning or ''}\nStudent sentence: {sentence}"
            raw = await _gemini_generate(GRADE_SENTENCE_SYSTEM, prompt, timeout=8.0)
            raw = raw.strip()
            if raw.startswith("```"):
                raw = re.sub(r"^```(?:json)?\s*", "", raw)
                raw = re.sub(r"\s*```$", "", raw)
            parsed = json.loads(raw)
            if isinstance(parsed, dict) and "passed" in parsed:
                result = {
                    "passed": bool(parsed.get("passed")),
                    "score_pct": int(parsed.get("score_pct") or (100 if parsed.get("passed") else 30)),
                    "feedback": str(parsed.get("feedback") or ""),
                    "source": "gemini",
                }
        except Exception as e:
            logger.debug("grade-sentence gemini fallback: %s", e)
    if not result:
        fb = _grade_sentence_fallback(kw, sentence, payload.meaning)
        result = {**fb, "source": "keyword"}
    result["passed"] = result.get("score_pct", 0) >= 60 and bool(result.get("passed"))
    return result


@api_router.post("/profile/avatar/unlock")
async def unlock_avatar_skin(payload: AvatarSkinRequest, owner: str = Depends(get_owner_id)):
    skin = next((s for s in AVATAR_CATALOG if s["skin_id"] == payload.skin_id), None)
    if not skin:
        raise HTTPException(status_code=404, detail="Unknown skin")
    profile = await get_device_profile(db, owner)
    unlocked = list(profile.get("unlocked_skins") or DEFAULT_UNLOCKED_SKINS)
    if payload.skin_id in unlocked:
        return {"unlocked_skins": unlocked, "diamonds": profile.get("diamonds", 0), "already_unlocked": True}
    diamonds = profile.get("diamonds", 0)
    if diamonds < skin["cost"]:
        raise HTTPException(status_code=402, detail=f"Not enough diamonds — need {skin['cost']}, have {diamonds}")
    unlocked.append(payload.skin_id)
    new_total = diamonds - skin["cost"]
    await db.device_profiles.update_one(
        {"owner": owner},
        {"$set": {"unlocked_skins": unlocked, "diamonds": new_total, "updated_at": datetime.now(timezone.utc).isoformat()}},
        upsert=True,
    )
    return {"unlocked_skins": unlocked, "diamonds": new_total, "already_unlocked": False}


@api_router.post("/profile/avatar/equip")
async def equip_avatar_skin(payload: AvatarSkinRequest, owner: str = Depends(get_owner_id)):
    profile = await get_device_profile(db, owner)
    unlocked = profile.get("unlocked_skins") or DEFAULT_UNLOCKED_SKINS
    if payload.skin_id not in unlocked:
        raise HTTPException(status_code=403, detail="Skin not unlocked yet")
    await db.device_profiles.update_one(
        {"owner": owner},
        {"$set": {"equipped_skin": payload.skin_id, "updated_at": datetime.now(timezone.utc).isoformat()}},
        upsert=True,
    )
    return {"equipped_skin": payload.skin_id}


# ============== SPACED REPETITION ENGINE (SM-2 + Forgetting Curve) ==============
MIN_EASE_FACTOR = 1.3
DEFAULT_EASE_FACTOR = 2.5
MEMORY_THRESHOLD = 70.0

RECOMMENDED_MIN_DAYS = {
    "reading_dictation": 3,
    "recital_dictation": 3,
    "quiz": 14,
    "exam": 28,
}


def _parse_dt(v):
    if not v:
        return None
    if isinstance(v, datetime):
        dt = v
    else:
        dt = datetime.fromisoformat(v)
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    return dt


def quality_from_response(correct: bool, response_time_ms: Optional[int] = None) -> int:
    if not correct:
        return 1
    if response_time_ms is None:
        return 4
    if response_time_ms < 3000:
        return 5
    if response_time_ms < 8000:
        return 4
    return 3


def sm2_update(ease_factor: float, interval_days: float, repetitions: int, quality: int):
    if quality < 3:
        repetitions = 0
        interval_days = 1
    else:
        if repetitions == 0:
            interval_days = 1
        elif repetitions == 1:
            interval_days = 6
        else:
            interval_days = round((interval_days or 1) * ease_factor)
        repetitions += 1

    ease_factor = ease_factor + (0.1 - (5 - quality) * (0.08 + (5 - quality) * 0.02))
    ease_factor = round(max(MIN_EASE_FACTOR, ease_factor), 3)
    return ease_factor, interval_days, repetitions


def compute_memory_strength(interval_days: float, last_reviewed_at: Optional[datetime], now: Optional[datetime] = None) -> float:
    if last_reviewed_at is None:
        return 0.0
    now = now or datetime.now(timezone.utc)
    last_reviewed_at = _parse_dt(last_reviewed_at)
    days_elapsed = (now - last_reviewed_at).total_seconds() / 86400.0
    if days_elapsed <= 0:
        return 100.0
    safe_interval = max(interval_days or 0.5, 0.5)
    strength = 100.0 * (0.9 ** (days_elapsed / safe_interval))
    return round(max(0.0, min(100.0, strength)), 1)


def _unit_memory_strength(u: dict, now: Optional[datetime] = None) -> float:
    return compute_memory_strength(u.get("interval_days", 0) or 0, _parse_dt(u.get("last_reviewed_at")), now)


def compute_lead_time_days(due_date: Optional[datetime], now: Optional[datetime] = None) -> Optional[int]:
    if not due_date:
        return None
    now = now or datetime.now(timezone.utc)
    due_date = _parse_dt(due_date)
    return max(0, (due_date - now).days)


# ============== URGENCY TIER (Section 10-11: Cram / Emergency / Survival) ==============
# Four tiers, computed LIVE from due_date vs. now (never a frozen flag), so a track
# quietly climbs the ladder as its deadline approaches even if nobody re-opens the
# app to "recompute" anything:
#   standard  -> lead time comfortably covers the Section 1 recommended minimum
#   cram      -> lead time is shorter than recommended, but still > 24h
#   emergency -> lead time <= 24h (Section 10.2/10.3 Emergency Journey)
#   survival  -> lead time <= 2h (Section 11 Survival/SOS Mode)
URGENCY_TIERS = ("standard", "cram", "emergency", "survival")

TIER_MAX_UNITS = {
    "standard": 8,
    "cram": 8,
    "emergency": 6,
    "survival": 5,
}

TIER_META = {
    "standard": {"label": "Standard Journey", "emoji": "🗺️"},
    "cram": {"label": "Cram Sprint", "emoji": "⚡"},
    "emergency": {"label": "Emergency Journey", "emoji": "🚨"},
    "survival": {"label": "SOS Mode", "emoji": "🆘"},
}


def compute_urgency_tier(track_type: str, due_date: Optional[datetime], now: Optional[datetime] = None) -> str:
    if not due_date:
        return "standard"
    now = now or datetime.now(timezone.utc)
    due_date = _parse_dt(due_date)
    lead_hours = (due_date - now).total_seconds() / 3600.0
    if lead_hours <= 2:
        return "survival"
    if lead_hours <= 24:
        return "emergency"
    min_days = RECOMMENDED_MIN_DAYS.get(track_type, 7)
    if (lead_hours / 24.0) < min_days:
        return "cram"
    return "standard"


def _attach_tier(track: dict, now: Optional[datetime] = None) -> dict:
    now = now or datetime.now(timezone.utc)
    tier = compute_urgency_tier(track.get("track_type"), track.get("due_date"), now)
    track["urgency_tier"] = tier
    track["urgency_meta"] = TIER_META[tier]
    track["is_cram"] = tier != "standard"  # kept for backward compatibility
    return track


def _select_units_for_tier(units: list, tier: str, max_units: Optional[int] = None) -> list:
    """Priority selection per tier. Section 10.1 Triage: in cram/emergency/survival,
    time goes to the weakest units first instead of spreading evenly across everything."""
    max_units = max_units or TIER_MAX_UNITS.get(tier, 8)
    due = sorted(
        [u for u in units if u.get("review_count", 0) > 0 and u.get("memory_strength", 0) < MEMORY_THRESHOLD],
        key=lambda u: u["memory_strength"],
    )
    new_units = [u for u in units if u.get("review_count", 0) == 0]
    selected = (due + new_units)[:max_units] or units[:max_units]
    return selected


def _apply_dictation_sampling(selected: list, track_type: str) -> list:
    """讀默 spec correction (Section B.2 / D.5): a real dictation only tests ~80%
    of the taught items, in shuffled order — never the full list in original
    order. Only meaningful once there's enough material to sample from."""
    if track_type != "reading_dictation" or len(selected) < 5:
        return selected
    sample_size = max(4, round(len(selected) * 0.8))
    sampled = random.sample(selected, min(sample_size, len(selected)))
    random.shuffle(sampled)
    return sampled


# ============== FOLLOW-UP REINFORCEMENT (Section 10.4 / 11.3 step S5) ==============
# Cram/Emergency/Survival readiness is, by design, unreliable — the whole point of
# those tiers is "pass tomorrow", not "remember long-term". So the moment such a
# session finishes, we quietly schedule a short retest a day or two later. If that
# retest still comes back weak, the units get flagged long_term_weak so future
# Weak Topic Drills (Section 4 Step 7 / Section 5 Step 5) pick them back up.
FOLLOWUP_DELAY_DAYS = {"cram": 2, "emergency": 1, "survival": 1}


async def _schedule_follow_up_if_needed(owner: str, track_id: Optional[str], unit_ids: List[str], now: datetime):
    if not track_id or not unit_ids:
        return
    track = await db.learning_tracks.find_one({"track_id": track_id, "owner": owner}, {"_id": 0})
    if not track:
        return
    tier = compute_urgency_tier(track.get("track_type"), track.get("due_date"), now)
    if tier == "standard":
        return
    # One pending follow-up per track at a time — don't stack duplicates if the
    # child plays several cram sessions on the same track before the retest is due.
    existing = await db.follow_up_tasks.find_one({"owner": owner, "track_id": track_id, "status": "pending"}, {"_id": 0})
    delay_days = FOLLOWUP_DELAY_DAYS.get(tier, 1)
    scheduled = now + timedelta(days=delay_days)
    if existing:
        merged_units = sorted(set(existing.get("unit_ids", [])) | set(unit_ids))
        await db.follow_up_tasks.update_one(
            {"follow_up_id": existing["follow_up_id"]},
            {"$set": {"unit_ids": merged_units, "tier": tier}},
        )
        return
    task = FollowUpTask(owner=owner, track_id=track_id, tier=tier, unit_ids=list(set(unit_ids)), scheduled_date=scheduled)
    doc = task.model_dump()
    doc["created_at"] = doc["created_at"].isoformat()
    doc["scheduled_date"] = scheduled.isoformat()
    await db.follow_up_tasks.insert_one(doc)


async def _maybe_resolve_follow_ups(owner: str, unit_id: str, now: datetime):
    """Called after every /reviews post. Only counts as 'the follow-up attempt' if
    it happens on/after the scheduled retest date — earlier reviews (e.g. more cram
    practice before the deadline) don't count, per Section 10.4's 1-2 day design."""
    pending = await db.follow_up_tasks.find(
        {"owner": owner, "status": "pending", "unit_ids": unit_id}, {"_id": 0}
    ).to_list(50)
    for fu in pending:
        sched = _parse_dt(fu.get("scheduled_date"))
        if sched and now < sched:
            continue
        reviewed = set(fu.get("reviewed_unit_ids", []))
        reviewed.add(unit_id)
        update = {"reviewed_unit_ids": list(reviewed)}
        if reviewed >= set(fu["unit_ids"]):
            fresh_units = await db.knowledge_units.find(
                {"unit_id": {"$in": fu["unit_ids"]}}, {"_id": 0}
            ).to_list(200)
            strengths = [_unit_memory_strength(u, now) for u in fresh_units]
            avg = sum(strengths) / len(strengths) if strengths else 0.0
            if avg < MEMORY_THRESHOLD:
                update["status"] = "escalated_to_weak_topic"
                await db.knowledge_units.update_many(
                    {"unit_id": {"$in": fu["unit_ids"]}}, {"$set": {"long_term_weak": True}}
                )
            else:
                update["status"] = "done"
        await db.follow_up_tasks.update_one({"follow_up_id": fu["follow_up_id"]}, {"$set": update})


def _shuffle(arr):
    import random
    a = list(arr)
    random.shuffle(a)
    return a


# ============== KNOWLEDGE UNIT EXTRACTION (LLM) ==============
UNIT_EXTRACTION_SYSTEM_PROMPT = """You are an expert K-12 curriculum analyst.
Break the given study material into discrete Knowledge Units for spaced-repetition learning.
Return ONLY valid JSON, no markdown fences, no commentary.

Schema:
{
  "units": [
    {"unit_type": "word", "term": "elephant", "meaning": "large description", "context": "optional sentence"}
  ]
}

Rules by track_type:
- reading_dictation: unit_type = "word". One unit per vocabulary word/character group.
- recital_dictation: unit_type = "sentence". One unit per sentence or short clause to memorize.
- quiz: unit_type = "concept". One unit per distinct concept/term being tested.
- exam: unit_type = "topic". One unit per topic/sub-topic.

Extract between 6 and 30 units depending on material length. Keep "term" short.
"meaning" should be a short definition/translation/explanation, or null.
"""


def _extract_units_fallback(text: str, track_type: str) -> list:
    # Pre-split lines (from Preview "Words / Sentences / Paragraphs" tool)
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    if lines and len(lines) >= 2:
        if track_type == "reading_dictation":
            return [{"unit_type": "word", "term": ln, "meaning": None, "context": None} for ln in lines[:40]]
        if track_type == "recital_dictation":
            return [{"unit_type": "sentence", "term": ln, "meaning": None, "context": None} for ln in lines[:30]]
        unit_type = "concept" if track_type == "quiz" else "topic"
        return [{"unit_type": unit_type, "term": ln[:120], "meaning": None, "context": None} for ln in lines[:20]]

    if track_type == "reading_dictation":
        words = [w for w in re.findall(r"[A-Za-z]+|[\u4e00-\u9fff]+", text) if len(w) > 1]
        seen, out = set(), []
        for w in words:
            if w not in seen:
                seen.add(w)
                out.append({"unit_type": "word", "term": w, "meaning": None, "context": None})
            if len(out) >= 20:
                break
        return out or [{"unit_type": "word", "term": "example", "meaning": None, "context": None}]
    elif track_type == "recital_dictation":
        sentences = [s.strip() for s in re.split(r"(?<=[。！？.!?])\s*", text) if len(s.strip()) > 3]
        return [{"unit_type": "sentence", "term": s, "meaning": None, "context": None} for s in sentences[:20]] or \
               [{"unit_type": "sentence", "term": text[:80], "meaning": None, "context": None}]
    else:
        unit_type = "concept" if track_type == "quiz" else "topic"
        chunks = [c.strip() for c in re.split(r"\n+", text) if len(c.strip()) > 5]
        return [{"unit_type": unit_type, "term": c[:60], "meaning": None, "context": None} for c in chunks[:15]] or \
               [{"unit_type": unit_type, "term": text[:60], "meaning": None, "context": None}]


async def _extract_knowledge_units(text: str, track_type: str) -> list:
    if not GEMINI_ENABLED:
        return _extract_units_fallback(text, track_type)
    try:
        prompt = f"track_type: {track_type}\n\nMaterial:\n---\n{text[:4000]}\n---\nReturn ONLY the JSON object."
        raw = await _gemini_generate(UNIT_EXTRACTION_SYSTEM_PROMPT, prompt)
        data = _extract_json(raw if isinstance(raw, str) else str(raw))
        units = data.get("units")
        if not units or not isinstance(units, list):
            raise ValueError("invalid units payload")
        return units
    except Exception as e:
        logging.exception("unit extraction failed, using fallback: %s", e)
        return _extract_units_fallback(text, track_type)


# ============== AI GAME GENERATION ==============
SYSTEM_PROMPT = """You are an expert Chinese & English K-12 learning game designer.
You convert student worksheets/textbook excerpts into a structured JSON game payload.
The game has 7 challenge types:
- "tap": tap the correct answer from 4 options (MCQ for vocab meaning, synonym, grammar, character recognition)
- "drag": drag word/character tiles to reconstruct a sentence in correct order
- "memory_flash": shows a word/phrase briefly, then user selects it from options
- "idiom_repair": show a Chinese idiom with one wrong/missing character; 4 single-character options; answer is the correct character
- "stroke_order": show a Chinese character; 4 options are stroke-sequence strings (e.g. "1-2-3-4" or next stroke number); answer is correct sequence
- "typing": real dictation simulation (讀默). The client plays the "answer" word/phrase
  aloud via text-to-speech and shows NO text or image hint at all — the student types
  what they hear. CRITICAL: the "prompt" field must be a generic instruction only
  (e.g. "Listen carefully and type what you hear") and must NEVER contain, hint at, or
  spell out the target word itself — the word only ever appears in the "answer" field.
- "full_recall": full-passage recitation check (背默). The student must type out an
  entire memorized passage from scratch with no visual reference. "prompt" must be a
  generic instruction (e.g. "背默：憑記憶寫出你啱啱背嘅內容" or "Write out the full
  passage from memory") and must NOT include or paraphrase the passage text. "answer"
  holds the full original passage text used for comparison. Only ever generate ONE
  "full_recall" challenge per game, covering the combined text of ALL recital_dictation
  units provided, and include ALL those units' unit_ids in an "answer_unit_ids" array
  on that challenge (the "unit_id" field on a full_recall challenge can be the first
  covered unit as a fallback).

You MUST return ONLY valid JSON. No markdown fences. No commentary.

Schema:
{
  "title": "Boss / Adventure title (short, exciting, in same language as content)",
  "boss_name": "Boss creature name fitting fantasy theme",
  "intro": "1-sentence mission briefing for the child",
  "challenges": [
    {
      "type": "tap",
      "prompt": "Question prompt",
      "options": ["A","B","C","D"],
      "answer": "correct option text",
      "explanation": "1 short sentence why"
    },
    {
      "type": "drag",
      "prompt": "Reconstruct this sentence",
      "tiles": ["I","love","reading","books"],
      "answer": ["I","love","reading","books"]
    },
    {
      "type": "memory_flash",
      "prompt": "Memorize then choose",
      "target": "word or phrase",
      "options": ["correct","distractor1","distractor2","distractor3"],
      "answer": "correct"
    },
    {
      "type": "idiom_repair",
      "prompt": "Fix the idiom",
      "idiom": "一__当先",
      "options": ["马","牛","羊","狗"],
      "answer": "马"
    },
    {
      "type": "stroke_order",
      "prompt": "Pick the correct stroke order",
      "character": "永",
      "options": ["1-2-3-4-5","2-1-3-4-5","1-3-2-4-5","1-2-4-3-5"],
      "answer": "1-2-3-4-5"
    },
    {
      "type": "typing",
      "prompt": "Listen carefully and type what you hear",
      "answer": "elephant"
    },
    {
      "type": "full_recall",
      "prompt": "背默：憑記憶寫出你啱啱背嘅內容",
      "answer": "the full original passage text goes here",
      "answer_unit_ids": ["ku_xxx", "ku_yyy"]
    }
  ]
}
"""


def _build_user_prompt(text: str, mode: str, subject: str, difficulty: int) -> str:
    mode_map = {
        "reading_dictation": (
            "讀默 (Reading Dictation): this MUST feel like a real school dictation, not a "
            "multiple-choice quiz — the real stimulus is the teacher's spoken voice, not "
            "text or pictures. At least 50% of challenges must be \"typing\" (audio-only, "
            "no visual hint, student types what they hear). The remaining challenges may "
            "be \"tap\"/\"memory_flash\" for the earlier recognition/understanding stage."
        ),
        "recital_dictation": (
            "背默 (Recital Dictation): focus on sentence reconstruction (drag), missing "
            "words, paragraph memory for most challenges. The pass bar is reproducing the "
            "FULL memorized passage, not fragments — you MUST include exactly one "
            "\"full_recall\" challenge whose answer is the complete passage text (or as "
            "much of it as fits), not just one sentence."
        ),
        "quiz": "Quiz: mixed challenges testing meaning, grammar, comprehension.",
        "exam": "Exam: harder mixed challenges with time pressure and tricky distractors.",
    }
    mode_desc = mode_map.get(mode, "Quiz")
    return f"""Subject: {subject}
Mode: {mode_desc}
Difficulty 1-5: {difficulty}

Source material:
---
{text[:3500]}
---

Generate EXACTLY 8 challenges total, well-balanced for the mode above.
- For Chinese content include at least 1 "idiom_repair" and 1 "stroke_order" when possible.
- For reading_dictation and recital_dictation, follow the per-mode "typing"/"full_recall"
  requirements above exactly — these override the general mix below.
- Otherwise (quiz/exam) ~40% "tap", ~25% "drag", ~20% "memory_flash", ~15% Chinese-specific types.
- Use the same language as the source (Chinese/English).
- Drag tile arrays must have between 3 and 7 tiles. For Chinese, tiles can be single characters or 2-char words.
- Options arrays must have exactly 4 items with one correct answer.
- "typing"/"full_recall" prompts must be generic instructions ONLY — never reveal, hint at, or paraphrase the answer text.
- Keep prompts short (under 60 chars), except full_recall's answer itself which may be long.
Return ONLY the JSON object.
"""


def _extract_json(s: str) -> dict:
    s = s.strip()
    if s.startswith("```"):
        s = re.sub(r"^```(?:json)?\s*", "", s)
        s = re.sub(r"\s*```$", "", s)
    # Find first { and last }
    start = s.find("{")
    end = s.rfind("}")
    if start != -1 and end != -1:
        s = s[start:end + 1]
    return json.loads(s)


def _has_cjk(text: str) -> bool:
    return bool(re.search(r"[\u4e00-\u9fff]", text))


def _fallback_game(text: str, mode: str) -> dict:
    # Build a simple deterministic game from the text in case LLM fails
    words = [w for w in re.findall(r"[A-Za-z]+|[\u4e00-\u9fff]+", text) if len(w) > 1][:16]
    if not words:
        words = ["hero", "magic", "study", "adventure", "learn", "brave", "battle", "win"]
    cjk = _has_cjk(text)
    challenges = []
    # 3 tap
    for i in range(3):
        target = words[i % len(words)]
        distractors = [w for w in words if w != target][:3] or ["alpha", "beta", "gamma"]
        opts = _shuffle((distractors + [target])[:4])
        challenges.append({
            "type": "tap",
            "prompt": "Which word appears in the material?",
            "options": opts,
            "answer": target,
            "explanation": "It comes from your uploaded text.",
        })
    if cjk:
        challenges.append({
            "type": "idiom_repair",
            "prompt": "Fix the idiom — pick the correct character",
            "idiom": "一__当先",
            "options": ["马", "牛", "羊", "狗"],
            "answer": "马",
        })
        challenges.append({
            "type": "stroke_order",
            "prompt": "Pick the correct stroke order",
            "character": words[0][0] if words[0] else "永",
            "options": ["1-2-3-4-5", "2-1-3-4-5", "1-3-2-4-5", "1-2-4-3-5"],
            "answer": "1-2-3-4-5",
        })
    # 2 drag
    for i in range(2):
        chunk = words[i * 3:(i * 3) + 4] or words[:4]
        challenges.append({
            "type": "drag",
            "prompt": "Arrange the tiles in order",
            "tiles": chunk,
            "answer": chunk,
        })
    # 2 memory flash
    for i in range(2):
        target = words[(i + 4) % len(words)]
        distractors = [w for w in words if w != target][:3] or ["x", "y", "z"]
        opts = _shuffle((distractors + [target])[:4])
        challenges.append({
            "type": "memory_flash",
            "prompt": "Memorize then pick",
            "target": target,
            "options": opts,
            "answer": target,
        })
    return {
        "title": "Quick Adventure",
        "boss_name": "Memory Goblin",
        "intro": "Defeat the goblin using words from your material!",
        "challenges": challenges[:8],
    }


@api_router.post("/generate-game")
async def generate_game(payload: GenerateGameRequest):
    if not GEMINI_ENABLED:
        return {"game": _fallback_game(payload.text, payload.mode), "source": "fallback"}

    user_prompt = _build_user_prompt(payload.text, payload.mode, payload.subject, payload.difficulty)
    try:
        raw = await _gemini_generate(SYSTEM_PROMPT, user_prompt)
        game = _extract_json(raw if isinstance(raw, str) else str(raw))
        if "challenges" not in game or not isinstance(game["challenges"], list) or len(game["challenges"]) == 0:
            raise ValueError("invalid game payload")
        return {"game": game, "source": "gemini"}
    except Exception as e:
        logging.exception("generate-game failed, using fallback: %s", e)
        return {"game": _fallback_game(payload.text, payload.mode), "source": "fallback"}


# ============== BATTLE FROM KNOWLEDGE UNITS (Daily Task Engine plugged into gameplay) ==============
BATTLE_FROM_UNITS_SYSTEM_PROMPT = """You are an expert K-12 learning game designer.
Given a list of Knowledge Units (each with a unit_id, term, and optional meaning/context),
create battle challenges that test EACH unit provided, one challenge per unit, same order.
Return ONLY valid JSON, no markdown fences, no commentary.

Each challenge type must be one of: "tap" (MCQ), "drag" (reconstruction), "memory_flash" (memorize
then recall), "typing" (audio-only dictation — no text/image hint, type what you hear), or
"full_recall" (type an entire passage from memory, no visual reference).
Pick whichever type best fits the unit content and the mode.

Mode-specific rules (check the "Mode:" line in the user message):
- Mode reading_dictation (讀默): the user message includes a "New units" list (unit_ids
  never reviewed before) and a "Reviewed units" list. For units in "New units", generate
  an easier recognition/understanding challenge ("tap" MCQ matching meaning, or
  "memory_flash") — NEVER "typing" for a unit the child has never seen tested before, that
  is developmentally unfair. For units in "Reviewed units", generate "typing" (audio-only
  dictation, no text/image hint). If "New units" is empty (everything has been reviewed at
  least once), then at least 50% of challenges should be "typing". "prompt" for typing
  challenges must be a generic instruction only (e.g. "Listen and type what you hear") and
  must NEVER contain, hint at, or spell out the unit's term — the term only belongs in
  "answer".
- Mode recital_dictation (背默): the user message includes the same New/Reviewed split at
  sentence level. Only include a "full_recall" (full blind passage recall) challenge if
  the passage's sentence units are ALL in "Reviewed units". If any sentence unit is still
  "New", generate "drag" (sentence reconstruction) or missing-word "tap" challenges for
  those instead, and skip full_recall this round — the child needs to encounter the
  passage at least once before being asked to reproduce all of it from memory blind.
  When you DO generate full_recall, combine ALL provided sentence units into a single
  full-passage "answer" (join their terms in the given order) and list every covered
  unit_id in "answer_unit_ids". "prompt" must be a generic instruction (e.g. "背默：憑記憶
  寫出你啱啱背嘅內容") — do NOT reveal or paraphrase the passage text in "prompt".

Schema:
{
  "title": "...", "boss_name": "...", "intro": "...",
  "challenges": [
    {"unit_id": "ku_xxx", "type": "tap", "prompt": "...", "options": ["A","B","C","D"], "answer": "...", "explanation": "..."},
    {"unit_id": "ku_xxx", "type": "drag", "prompt": "...", "tiles": ["..."], "answer": ["..."]},
    {"unit_id": "ku_xxx", "type": "memory_flash", "prompt": "...", "target": "...", "options": ["..."], "answer": "..."},
    {"unit_id": "ku_xxx", "type": "typing", "prompt": "Listen and type what you hear", "answer": "..."},
    {"unit_id": "ku_xxx", "type": "full_recall", "prompt": "背默：憑記憶寫出你啱啱背嘅內容", "answer": "...", "answer_unit_ids": ["ku_xxx","ku_yyy"]}
  ]
}
Every challenge MUST include the correct "unit_id" copied verbatim from the input list
(for full_recall, "unit_id" can be the first covered unit as a fallback reference).
"""


def _fallback_game_from_units(units: list, track_type: Optional[str] = None, new_ids: Optional[list] = None) -> dict:
    terms = [u["term"] for u in units] or ["example"]
    new_id_set = set(new_ids or [])

    if track_type == "reading_dictation" and units:
        # B.2 correction + first-exposure fairness: units never reviewed before get an
        # easier recognition MCQ instead of blind audio typing — only units the child has
        # already been through at least once get the real dictation simulation.
        challenges = []
        for u in units:
            if u["unit_id"] in new_id_set:
                target = u["term"]
                distractors = [t for t in terms if t != target][:3] or ["alpha", "beta", "gamma"]
                challenges.append({
                    "unit_id": u["unit_id"],
                    "type": "tap",
                    "prompt": f"Which one matches: {u.get('meaning') or 'this word'}?",
                    "options": _shuffle((distractors + [target])[:4]),
                    "answer": target,
                    "explanation": u.get("meaning") or "First time seeing this one — get familiar with it!",
                })
            else:
                challenges.append({
                    "unit_id": u["unit_id"],
                    "type": "typing",
                    "prompt": "🔊 聽讀音，寫出正確答案 / Listen and type what you hear",
                    "answer": u["term"],
                    "explanation": u.get("meaning") or "From your material.",
                })
        return {
            "title": "Dictation Trial",
            "boss_name": "Echo Wraith",
            "intro": "Listen carefully — the Echo Wraith only fears words spelled from memory!",
            "challenges": challenges,
        }

    if track_type == "recital_dictation" and units:
        # C correction: 背默 pass bar is full-passage reproduction — but only once every
        # sentence has been seen at least once. If any sentence is brand new, warm up with
        # reconstruction/missing-word practice first instead of a blind full-passage test.
        if new_id_set & {u["unit_id"] for u in units}:
            challenges = []
            for u in units:
                words = u["term"].split() if " " in u["term"] else list(u["term"])
                challenges.append({
                    "unit_id": u["unit_id"],
                    "type": "drag",
                    "prompt": "重組句子 / Arrange in order",
                    "tiles": _shuffle(words) if len(words) > 1 else words,
                    "answer": words,
                })
            return {
                "title": "Passage Warm-up",
                "boss_name": "Scroll Apprentice",
                "intro": "Get familiar with the passage first — piece it back together!",
                "challenges": challenges,
            }
        full_text = "".join(u["term"] for u in units) if any(len(u["term"]) <= 2 for u in units) \
            else " ".join(u["term"] for u in units)
        challenges = [{
            "unit_id": units[0]["unit_id"],
            "type": "full_recall",
            "prompt": "背默：憑記憶寫出你啱啱背嘅內容 / Write out the full passage from memory",
            "answer": full_text,
            "answer_unit_ids": [u["unit_id"] for u in units],
        }]
        return {
            "title": "Full Recital Trial",
            "boss_name": "Scroll Keeper",
            "intro": "Recite the whole passage from memory — no peeking, no fragments!",
            "challenges": challenges,
        }

    challenges = []
    for u in units:
        target = u["term"]
        distractors = [t for t in terms if t != target][:3] or ["alpha", "beta", "gamma"]
        opts = _shuffle((distractors + [target])[:4])
        challenges.append({
            "unit_id": u["unit_id"],
            "type": "tap",
            "prompt": f"Which is correct: {u.get('meaning') or 'pick the matching item'}?",
            "options": opts,
            "answer": target,
            "explanation": u.get("meaning") or "From your material.",
        })
    return {
        "title": "Quick Adventure",
        "boss_name": "Memory Goblin",
        "intro": "Defeat the goblin using your knowledge units!",
        "challenges": challenges,
    }


# ============== SESSION-INTERNAL PACING (Section 10.1 Expanding Retrieval / 11.3 Rapid Loop) ==============
# These two tiers don't have the luxury of real 24-72hr spaced intervals, so instead
# of a flat "one challenge per unit, once" pass, we reshape the SAME set of
# challenges into a sequence that manufactures artificial spacing (emergency) or
# leans into massed repetition on purpose (survival) — both are legitimate
# techniques for compressed timelines, not a lesser version of the real thing.
EMERGENT_CHUNK_SIZE = 3
SURVIVAL_REPEATS = 2  # Section 11.3 S3: "重複2-3次"


def _expanding_retrieval_order(unit_ids: List[str]) -> List[str]:
    """Section 10.1/10.2 (E2-E4): learn a small chunk, then while learning the NEXT
    chunk, force a quick flashback recall of every earlier chunk — the interval
    between first-seeing an item and being asked again keeps expanding as more
    chunks get added, mimicking spaced repetition inside a single sitting."""
    chunks = [unit_ids[i:i + EMERGENT_CHUNK_SIZE] for i in range(0, len(unit_ids), EMERGENT_CHUNK_SIZE)]
    seq: List[str] = []
    seen_chunks: List[List[str]] = []
    for chunk in chunks:
        seq.extend(chunk)  # first exposure to this chunk
        for prev in seen_chunks:
            seq.extend(prev)  # flashback retrieval of every earlier chunk, interleaved
        seen_chunks.append(chunk)
    # E5: one final full mixed-order recall round, per Section 10.2.
    tail = list(unit_ids)
    random.shuffle(tail)
    seq.extend(tail)
    return seq


def _rapid_loop_order(unit_ids: List[str]) -> List[str]:
    """Section 11.3 S3-S4: massed practice — show/hide/recall the SAME item 2-3
    times back-to-back (no spacing; the point is maximum reps in minimum time),
    then one final mixed round (S4) simulating the real dictation/test."""
    seq: List[str] = []
    for uid in unit_ids:
        seq.extend([uid] * SURVIVAL_REPEATS)
    tail = list(unit_ids)
    random.shuffle(tail)
    seq.extend(tail)
    return seq


def _sequence_challenges_for_tier(challenges: List[dict], unit_ids: List[str], tier: str) -> List[dict]:
    by_unit = {c.get("unit_id"): c for c in challenges if c.get("unit_id")}
    ordered_ids = [uid for uid in unit_ids if uid in by_unit]
    if not ordered_ids:
        return challenges

    order: List[str]
    if tier == "emergency":
        order = _expanding_retrieval_order(ordered_ids)
    elif tier == "survival":
        order = _rapid_loop_order(ordered_ids)
    else:
        return challenges  # standard/cram: single pass, order already Triage-prioritised on selection

    # Re-use the same challenge payload each time a unit repeats — BattlePage records
    # a /reviews call per occurrence, which is exactly the point (repetition IS the practice).
    return [dict(by_unit[uid]) for uid in order]


async def _generate_battle_from_units(selected: list, unit_ids: list, track_type: str, tier: str) -> dict:
    """Shared LLM-or-fallback game generation, with tier-aware framing per Section 10.5/11.4:
    cram/emergency/survival get an urgent, action-oriented tone instead of the calm standard one."""
    new_ids = [u["unit_id"] for u in selected if u.get("review_count", 0) == 0]
    reviewed_ids = [u["unit_id"] for u in selected if u.get("review_count", 0) > 0]
    if not GEMINI_ENABLED:
        game = _fallback_game_from_units(selected, track_type, new_ids)
    else:
        try:
            tier_note = {
                "standard": "Tone: calm adventurous exploration.",
                "cram": "Tone: focused sprint training camp — energetic but not scary. Call it a 'training sprint'.",
                "emergency": "Tone: urgent but encouraging special mission ('Emergency Mission'). Never say the child is late or behind.",
                "survival": "Tone: SOS special-ops mission ('SOS Mission'). Upbeat, action-focused, never mention lack of time or being behind schedule.",
            }.get(tier, "")
            unit_payload = [
                {"unit_id": u["unit_id"], "term": u["term"], "meaning": u.get("meaning"), "context": u.get("context")}
                for u in selected
            ]
            prompt = (
                f"Mode: {track_type}\n{tier_note}\n"
                f"New units (never reviewed — use recognition challenges, NOT typing/full_recall): {json.dumps(new_ids)}\n"
                f"Reviewed units (safe for typing/full_recall): {json.dumps(reviewed_ids)}\n"
                f"Knowledge Units:\n{json.dumps(unit_payload, ensure_ascii=False)}\n"
                "Return ONLY the JSON object."
            )
            raw = await _gemini_generate(BATTLE_FROM_UNITS_SYSTEM_PROMPT, prompt)
            game = _extract_json(raw if isinstance(raw, str) else str(raw))
            returned_ids = set()
            for c in game.get("challenges", []):
                if c.get("unit_id"):
                    returned_ids.add(c["unit_id"])
                for uid in c.get("answer_unit_ids") or []:
                    returned_ids.add(uid)
            if not game.get("challenges") or not returned_ids.issubset(set(unit_ids)):
                raise ValueError("invalid or mismatched unit_ids in generated game")
            # Safety net: even if the LLM ignores the instruction, never let a truly new
            # unit be tested with blind typing/full_recall — downgrade tap-able cases
            # deterministically, or fall back entirely if a full_recall covers new units
            # (can't cleanly downgrade "reproduce the whole passage" into anything fair).
            new_id_set = set(new_ids)
            needs_fallback = False
            for c in game["challenges"]:
                covered = set(c.get("answer_unit_ids") or ([c["unit_id"]] if c.get("unit_id") else []))
                if c.get("type") == "full_recall" and covered & new_id_set:
                    needs_fallback = True
                    break
                if c.get("type") == "typing" and covered & new_id_set:
                    c["type"] = "tap"
                    term = c.get("answer", "")
                    c["options"] = _shuffle([term, "...", "...", "..."])
            if needs_fallback:
                game = _fallback_game_from_units(selected, track_type, new_ids)
        except Exception as e:
            logging.exception("battle-from-units generation failed, using fallback: %s", e)
            game = _fallback_game_from_units(selected, track_type, new_ids)

    if game.get("challenges"):
        game["challenges"] = _sequence_challenges_for_tier(game["challenges"], unit_ids, tier)
        # A flat 100 HP assumed ~8 challenges. Expanding Retrieval / Rapid Loop can
        # produce far more (repetition IS the point) — without this, the boss dies
        # and the session ends before the planned repetitions actually happen.
        game["boss_max_hp"] = max(120, round(len(game["challenges"]) * 18))
    return game


@api_router.post("/tracks/{track_id}/battle")
async def start_track_battle(track_id: str, owner: str = Depends(get_owner_id), max_units: Optional[int] = None):
    track = await db.learning_tracks.find_one({"track_id": track_id, **_track_access_filter(owner)}, {"_id": 0})
    if not track:
        raise HTTPException(status_code=404, detail="Track not found")

    units = await db.knowledge_units.find({"track_id": track_id}, {"_id": 0}).to_list(500)
    if not units:
        raise HTTPException(status_code=400, detail="No knowledge units in this track")
    now = datetime.now(timezone.utc)
    for u in units:
        u["memory_strength"] = _unit_memory_strength(u, now)

    tier = compute_urgency_tier(track.get("track_type"), track.get("due_date"), now)
    selected = _select_units_for_tier(units, tier, max_units)
    selected = _apply_dictation_sampling(selected, track["track_type"])
    unit_ids = [u["unit_id"] for u in selected]

    game = await _generate_battle_from_units(selected, unit_ids, track["track_type"], tier)

    return {
        "game": game,
        "unit_ids": unit_ids,
        "track": {
            "track_id": track_id,
            "is_cram": tier != "standard",
            "urgency_tier": tier,
            "urgency_meta": TIER_META[tier],
        },
    }


@api_router.get("/daily-battle")
async def get_daily_battle(owner: str = Depends(get_owner_id), max_units: int = 8):
    """Cross-track Daily Task Engine battle (Section 7.2): pulls the same due+new
    units as /daily-queue, prioritising the most urgent track's units first, and
    turns them into a single playable battle instead of just a list."""
    tracks = await db.learning_tracks.find(_student_tracks_filter(owner), {"_id": 0}).to_list(200)
    if not tracks:
        raise HTTPException(status_code=400, detail="No active learning tracks")
    now = datetime.now(timezone.utc)
    track_map = {}
    for t in tracks:
        tier = compute_urgency_tier(t.get("track_type"), t.get("due_date"), now)
        track_map[t["track_id"]] = {**t, "urgency_tier": tier}

    units = await db.knowledge_units.find({"track_id": {"$in": list(track_map.keys())}}, {"_id": 0}).to_list(2000)
    followup_due_ids = set()
    followups = await db.follow_up_tasks.find({"owner": owner, "status": "pending"}, {"_id": 0}).to_list(200)
    for f in followups:
        sched = _parse_dt(f.get("scheduled_date"))
        if sched and sched <= now:
            followup_due_ids.update(f.get("unit_ids", []))
    for u in units:
        u["memory_strength"] = _unit_memory_strength(u, now)
        u["urgency_tier"] = track_map.get(u["track_id"], {}).get("urgency_tier", "standard")
        u["is_followup"] = u["unit_id"] in followup_due_ids

    tier_rank = {"survival": 0, "emergency": 1, "cram": 2, "standard": 3}

    def _queue_key(u):
        tr = track_map.get(u["track_id"], {})
        sp = 1 if tr.get("track_type") == "self_practice" else 0
        return (
            0 if u["is_followup"] else 1,
            sp,
            tier_rank.get(u["urgency_tier"], 3),
            u["memory_strength"],
        )

    due = sorted(
        [u for u in units if (u.get("review_count", 0) > 0 and u["memory_strength"] < MEMORY_THRESHOLD) or u["is_followup"]],
        key=_queue_key,
    )
    new_units = sorted(
        [u for u in units if u.get("review_count", 0) == 0],
        key=lambda u: _queue_key(u)[1:],
    )
    selected = (due + new_units)[:max_units]
    if not selected:
        raise HTTPException(status_code=400, detail="Nothing due today")
    unit_ids = [u["unit_id"] for u in selected]

    # Most urgent tier present drives the framing of this mixed battle.
    overall_tier = min((u["urgency_tier"] for u in selected), key=lambda t: tier_rank.get(t, 3))
    dominant_track_type = selected[0].get("track_type", "quiz")
    if all(u.get("track_type") == "reading_dictation" for u in selected):
        selected = _apply_dictation_sampling(selected, "reading_dictation")
        unit_ids = [u["unit_id"] for u in selected]

    game = await _generate_battle_from_units(selected, unit_ids, dominant_track_type, overall_tier)

    return {
        "game": game,
        "unit_ids": unit_ids,
        "urgency_tier": overall_tier,
        "urgency_meta": TIER_META[overall_tier],
    }


# ============== STREAK REWARDS ==============
async def _update_streak(owner: str, now: datetime) -> dict:
    today = now.date().isoformat()
    doc = await db.streaks.find_one({"owner": owner}, {"_id": 0})
    if not doc:
        doc = {"owner": owner, "current_streak": 1, "longest_streak": 1, "last_active_date": today}
        await db.streaks.insert_one({**doc})
        return doc
    last = doc.get("last_active_date")
    if last == today:
        return doc
    yesterday = (now.date() - timedelta(days=1)).isoformat()
    if last == yesterday:
        current = doc.get("current_streak", 0) + 1
    else:
        current = 1
    longest = max(doc.get("longest_streak", 0), current)
    doc = {"owner": owner, "current_streak": current, "longest_streak": longest, "last_active_date": today}
    await db.streaks.update_one({"owner": owner}, {"$set": doc}, upsert=True)
    return doc


@api_router.get("/streak")
async def get_streak(owner: str = Depends(get_owner_id)):
    doc = await db.streaks.find_one({"owner": owner}, {"_id": 0})
    if not doc:
        return {"owner": owner, "current_streak": 0, "longest_streak": 0, "last_active_date": None}
    return doc


# ============== PUSH NOTIFICATIONS (Capacitor native shell) ==============
# Registers an FCM (Android) / APNs (iOS) device token obtained client-side via
# @capacitor/push-notifications. We only store tokens here — actually *sending*
# pushes (e.g. daily "3 words due today" reminders) needs a Firebase Admin /
# APNs server integration, which is a separate follow-up (see FOREVER_NOTES.md).
class PushTokenRegister(BaseModel):
    token: str
    platform: str  # "ios" | "android"


@api_router.post("/push/register-token")
async def register_push_token(payload: PushTokenRegister, owner: str = Depends(get_owner_id)):
    if not payload.token:
        raise HTTPException(status_code=400, detail="token required")
    doc = {
        "owner": owner,
        "token": payload.token,
        "platform": payload.platform,
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.push_tokens.update_one(
        {"owner": owner, "token": payload.token}, {"$set": doc}, upsert=True
    )
    return {"ok": True}


@api_router.delete("/push/register-token")
async def unregister_push_token(token: str, owner: str = Depends(get_owner_id)):
    await db.push_tokens.delete_one({"owner": owner, "token": token})
    return {"ok": True}


PUSH_REMINDER_COOLDOWN_HOURS = int(os.environ.get("PUSH_REMINDER_COOLDOWN_HOURS", "6"))

PUSH_REMINDER_MESSAGES = [
    "{n}個字準備緊同你講拜拜，返嚟救返佢哋！",
    "你嘅背默進度就快跌落安全線，5分鐘幫佢充充電！",
    "記憶種子需要澆水啦 — {n}個字等你複習！",
    "🐉 限時 Boss 出現！{n}個字快忘記咗 — 返嚟救佢哋！",
]

GRADUATION_BOSS_DELAY_DAYS = 2


async def _schedule_graduation_boss_followup(track: dict, now: datetime):
    """After final step pass — schedule forgetting-curve timed boss (replay track anytime)."""
    owner = track.get("owner")
    track_id = track.get("track_id")
    if not owner or not track_id:
        return
    existing = await db.follow_up_tasks.find_one(
        {"owner": owner, "track_id": track_id, "kind": "graduation_boss", "status": "pending"},
        {"_id": 0},
    )
    if existing:
        return
    units = await db.knowledge_units.find({"track_id": track_id}, {"_id": 0}).to_list(500)
    unit_ids = [u["unit_id"] for u in units if u.get("unit_id")]
    if not unit_ids:
        return
    scheduled = now + timedelta(days=GRADUATION_BOSS_DELAY_DAYS)
    task = FollowUpTask(
        owner=owner,
        track_id=track_id,
        tier="standard",
        unit_ids=unit_ids,
        scheduled_date=scheduled,
    )
    doc = task.model_dump()
    doc["kind"] = "graduation_boss"
    doc["created_at"] = doc["created_at"].isoformat()
    doc["scheduled_date"] = scheduled.isoformat()
    await db.follow_up_tasks.insert_one(doc)


@api_router.get("/push/due-reminder")
async def push_due_reminder(owner: str = Depends(get_owner_id)):
    """Forgetting-curve reminder payload (v3 §6.1). Also sends FCM when configured."""
    payload = await _build_due_reminder_for_owner(owner)
    if not payload.get("should_notify"):
        return payload
    dl = payload["deep_link"]
    fcm_result = await send_fcm_to_owner(
        db,
        owner,
        title=payload["title"],
        body=payload["message"],
        data={
            "type": dl.get("type", "journey"),
            "track_id": dl.get("track_id", ""),
            "step": str(dl.get("step", "")),
            "bundle_index": str(dl.get("bundle_index", "")),
            "track_type": dl.get("track_type", ""),
        },
    )
    return {**payload, "fcm": fcm_result}


async def _build_due_reminder_for_owner(owner: str) -> dict:
    tracks = await db.learning_tracks.find(_student_tracks_filter(owner), {"_id": 0}).to_list(200)
    if not tracks:
        return {"should_notify": False, "due_count": 0}

    now = datetime.now(timezone.utc)
    track_map = {t["track_id"]: t for t in tracks}
    units = await db.knowledge_units.find(
        {"track_id": {"$in": list(track_map.keys())}}, {"_id": 0}
    ).to_list(2000)

    due_by_track: dict = {}
    for u in units:
        strength = _unit_memory_strength(u, now)
        if u.get("review_count", 0) > 0 and strength < MEMORY_THRESHOLD:
            tid = u["track_id"]
            due_by_track.setdefault(tid, []).append({**u, "memory_strength": strength})

    total_due = sum(len(v) for v in due_by_track.values())
    if total_due == 0:
        return {"should_notify": False, "due_count": 0}

    profile = await get_device_profile(db, owner)
    last_sent = _parse_dt(profile.get("last_push_reminder_at"))
    if last_sent and (now - last_sent).total_seconds() < PUSH_REMINDER_COOLDOWN_HOURS * 3600:
        return {"should_notify": False, "due_count": total_due, "cooldown": True}

    primary_tid = max(due_by_track.keys(), key=lambda tid: len(due_by_track[tid]))
    primary_track = track_map[primary_tid]
    bundle_idx = int(primary_track.get("current_bundle_index") or 0)
    step = int(primary_track.get("current_step") or 1)
    n = len(due_by_track[primary_tid])
    msg = random.choice(PUSH_REMINDER_MESSAGES).format(n=n)
    deep_link = {
        "type": "journey",
        "track_id": primary_tid,
        "step": step,
        "bundle_index": bundle_idx,
        "track_type": primary_track.get("track_type"),
    }
    await db.device_profiles.update_one(
        {"owner": owner},
        {"$set": {"last_push_reminder_at": now.isoformat()}},
        upsert=True,
    )
    return {
        "should_notify": True,
        "due_count": total_due,
        "message": msg,
        "title": "特工召回令 🔔",
        "deep_link": deep_link,
    }


@api_router.post("/push/process-due-reminders")
async def process_due_reminders_cron(
    x_cron_secret: Optional[str] = Header(default=None, alias="X-Cron-Secret"),
):
    """Batch send forgetting-curve FCM reminders. Set PUSH_CRON_SECRET in env."""
    expected = os.environ.get("PUSH_CRON_SECRET", "").strip()
    if not expected or x_cron_secret != expected:
        raise HTTPException(status_code=403, detail="Forbidden")

    owners = await db.push_tokens.distinct("owner")
    processed, notified = 0, 0
    for owner in owners:
        processed += 1
        payload = await _build_due_reminder_for_owner(owner)
        if not payload.get("should_notify"):
            continue
        dl = payload["deep_link"]
        await send_fcm_to_owner(
            db,
            owner,
            title=payload["title"],
            body=payload["message"],
            data={
                "type": dl.get("type", "journey"),
                "track_id": dl.get("track_id", ""),
                "step": str(dl.get("step", "")),
                "bundle_index": str(dl.get("bundle_index", "")),
                "track_type": dl.get("track_type", ""),
            },
        )
        notified += 1
    return {"processed": processed, "notified": notified, "fcm_configured": fcm_configured()}


@api_router.post("/behavioral-signals")
async def create_behavioral_signal(payload: BehavioralSignalCreate, owner: str = Depends(get_owner_id)):
    doc = await record_behavioral_signal(
        db,
        owner,
        unit_id=payload.unit_id,
        track_id=payload.track_id,
        game_type=payload.game_type,
        reaction_time_ms=payload.reaction_time_ms,
        processing_time_ms=payload.processing_time_ms,
        replay_count=payload.replay_count,
        hint_usage=payload.hint_usage,
        correct=payload.correct,
    )
    if not doc:
        return {"recorded": False, "reason": "consent_required", "consent_type": BEHAVIORAL_CONSENT_TYPE}
    return {"recorded": True, "signal": doc}


@api_router.get("/parent/behavioral-observations")
async def parent_behavioral_observations(
    kid_owner_id: str,
    user: Optional[User] = Depends(get_current_user),
):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    link = await db.family_links.find_one(
        {"parent_user_id": user.user_id, "kid_owner_id": kid_owner_id, "status": "active"},
        {"_id": 0},
    )
    if not link and kid_owner_id != user.user_id:
        raise HTTPException(status_code=403, detail="Not linked to this child")
    return await summarize_for_parent(db, kid_owner_id)


# ============== EMAIL PREFERENCES ==============
class EmailPreferences(BaseModel):
    enabled: bool = False
    frequency: str = "weekly"  # daily | weekly
    email: Optional[str] = None


@api_router.get("/follow-ups/email-preferences")
async def get_email_preferences(user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    doc = await db.email_preferences.find_one({"user_id": user.user_id}, {"_id": 0})
    if not doc:
        return {"enabled": False, "frequency": "weekly", "email": user.email}
    return doc


@api_router.post("/follow-ups/email-preferences")
async def set_email_preferences(payload: EmailPreferences, user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    doc = {
        "user_id": user.user_id,
        "enabled": payload.enabled,
        "frequency": payload.frequency,
        "email": payload.email or user.email,
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.email_preferences.update_one({"user_id": user.user_id}, {"$set": doc}, upsert=True)
    return doc


async def _maybe_send_progress_digest(user_id: str, owner: str):
    prefs = await db.email_preferences.find_one({"user_id": user_id, "enabled": True}, {"_id": 0})
    if not prefs or not prefs.get("email"):
        return
    sessions = await db.game_sessions.find({"owner": owner}, {"_id": 0}).to_list(100)
    total_score = sum(s.get("score", 0) for s in sessions)
    c = sum(s.get("correct", 0) for s in sessions)
    w = sum(s.get("wrong", 0) for s in sessions)
    streak_doc = await db.streaks.find_one({"owner": owner}, {"_id": 0})
    await send_progress_digest(
        prefs["email"],
        {
            "total_sessions": len(sessions),
            "total_score": total_score,
            "accuracy_pct": round(100 * c / (c + w)) if c + w else 0,
            "current_streak": streak_doc.get("current_streak", 0) if streak_doc else 0,
        },
        prefs.get("frequency", "weekly"),
    )


# ============== GAME SESSIONS ==============
DIAMONDS_PER_BOSS_DEFEAT = int(os.environ.get("DIAMONDS_PER_BOSS_DEFEAT", "5"))
DIAMONDS_PER_N_CORRECT = int(os.environ.get("DIAMONDS_PER_N_CORRECT", "5"))  # 1 diamond per N correct answers


DIAMONDS_PER_N_CORRECT = int(os.environ.get("DIAMONDS_PER_N_CORRECT", "5"))  # 1 diamond per N correct answers


async def _fetch_game_mode_performance(owner: str) -> dict:
    """Load success rates for game_mode_performance weighting (v3 §6.3a)."""
    docs = await db.game_mode_performance.find({"student_id": owner}, {"_id": 0}).to_list(50)
    out = {}
    for d in docs:
        gt = d.get("game_type")
        if not gt:
            continue
        attempts = int(d.get("attempts") or 0)
        out[gt] = {
            "attempts": attempts,
            "success_rate": (int(d.get("correct_count") or 0) / attempts) if attempts else 0,
        }
    return out


async def _record_game_mode_performance(owner: str, breakdown: Optional[dict]) -> None:
    if not breakdown:
        return
    now = datetime.now(timezone.utc).isoformat()
    for game_type, stats in breakdown.items():
        if not game_type or not isinstance(stats, dict):
            continue
        c = int(stats.get("correct") or 0)
        w = int(stats.get("wrong") or 0)
        n = c + w
        if n <= 0:
            continue
        existing = await db.game_mode_performance.find_one(
            {"student_id": owner, "game_type": game_type}, {"_id": 0}
        )
        if existing:
            attempts = int(existing.get("attempts") or 0) + n
            correct_count = int(existing.get("correct_count") or 0) + c
        else:
            attempts = n
            correct_count = c
        await db.game_mode_performance.update_one(
            {"student_id": owner, "game_type": game_type},
            {
                "$set": {
                    "attempts": attempts,
                    "correct_count": correct_count,
                    "success_rate": round(correct_count / attempts, 4) if attempts else 0,
                    "last_updated": now,
                }
            },
            upsert=True,
        )


@api_router.post("/game-sessions")
async def save_game_session(
    payload: GameSessionCreate,
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
):
    await assert_can_collect_learning_data(db, owner)
    now = datetime.now(timezone.utc)
    sess = GameSession(owner=owner, **payload.model_dump())
    doc = sess.model_dump()
    doc["completed_at"] = doc["completed_at"].isoformat()
    await db.game_sessions.insert_one({**doc})
    streak = await _update_streak(owner, now)
    if payload.track_id and payload.unit_ids:
        await _schedule_follow_up_if_needed(owner, payload.track_id, payload.unit_ids, now)
    if payload.track_id and payload.journey_step:
        journey_result = await _advance_journey_step(
            payload.track_id, payload.journey_step, payload.correct, payload.wrong
        )
    else:
        journey_result = {}
    await _record_game_mode_performance(owner, payload.game_breakdown)
    if user:
        await _maybe_send_progress_digest(user.user_id, owner)

    # Diamonds are a separate currency from score — Section 3 (badge/reward spec):
    # spent on cosmetic avatar skins/bonuses, never affects gameplay difficulty.
    diamonds_earned = (payload.correct // DIAMONDS_PER_N_CORRECT) + (
        DIAMONDS_PER_BOSS_DEFEAT if payload.boss_defeated else 0
    )
    coins_earned = payload.correct * COINS_PER_CORRECT
    diamonds_total = None
    coins_total = None
    inc_fields: dict = {}
    if diamonds_earned > 0:
        inc_fields["diamonds"] = diamonds_earned
    if coins_earned > 0:
        inc_fields["coins"] = coins_earned
    if inc_fields:
        updated = await db.device_profiles.find_one_and_update(
            {"owner": owner},
            {"$inc": inc_fields, "$setOnInsert": {"subscription_tier": "free"}},
            upsert=True,
            return_document=True,
        )
        diamonds_total = (updated or {}).get("diamonds", diamonds_earned)
        coins_total = (updated or {}).get("coins", coins_earned)

    session_count = await db.game_sessions.count_documents({"owner": owner})
    track_completed = journey_result.get("status") == "completed"
    new_badges = await check_and_award_badges(
        db,
        owner,
        session_count=session_count,
        streak=int((streak or {}).get("current_streak") or 0),
        max_combo=payload.max_combo,
        track_completed=track_completed,
    )

    out = {
        **doc,
        "streak": streak,
        "diamonds_earned": diamonds_earned,
        "diamonds_total": diamonds_total,
        "coins_earned": coins_earned,
        "coins_total": coins_total,
        "new_badges": new_badges,
        "journey": journey_result,
    }
    await db.progress_snapshots.delete_one({"owner": owner})
    return out


@api_router.get("/game-sessions")
async def list_game_sessions(
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
    kid_owner_id: Optional[str] = None,
):
    data_owner = await _resolve_child_data_owner(user, owner, kid_owner_id)
    if data_owner is None:
        return []
    docs = await db.game_sessions.find({"owner": data_owner}, {"_id": 0}).sort("completed_at", -1).to_list(100)
    return docs


# ============== LEARNING TRACKS (knowledge_unit + Memory Strength) ==============
@api_router.post("/tracks")
async def create_track(
    payload: TrackCreate,
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
):
    student_id, assigned_by = await _resolve_track_student(owner, user, payload)
    unit_owner = student_id or owner
    await assert_can_collect_learning_data(db, unit_owner)
    material_id = payload.material_id
    text = payload.text

    if material_id:
        mat_doc = await db.materials.find_one({"material_id": material_id}, {"_id": 0})
        if not mat_doc:
            raise HTTPException(status_code=404, detail="Material not found")
        text = mat_doc["text"]
    elif text:
        mat = Material(owner=owner, title=payload.title or "Untitled Material", text=text)
        doc = mat.model_dump()
        doc["created_at"] = doc["created_at"].isoformat()
        await db.materials.insert_one(doc)
        material_id = mat.material_id
    else:
        raise HTTPException(status_code=400, detail="Provide material_id or text")

    due_date_dt = None
    if payload.due_date and payload.track_type != "self_practice":
        try:
            due_date_dt = datetime.fromisoformat(payload.due_date.replace("Z", "+00:00"))
            if due_date_dt.tzinfo is None:
                due_date_dt = due_date_dt.replace(tzinfo=timezone.utc)
        except ValueError:
            due_date_dt = None

    lead_time_days = compute_lead_time_days(due_date_dt)
    now0 = datetime.now(timezone.utc)
    urgency_tier = compute_urgency_tier(payload.track_type, due_date_dt, now0)
    is_cram = urgency_tier != "standard" and payload.track_type != "self_practice"

    track = LearningTrack(
        owner=owner if not assigned_by else (assigned_by if payload.is_template else owner),
        student_id=student_id,
        assigned_by=assigned_by,
        material_id=material_id,
        track_type=payload.track_type,
        scope_description=(payload.title or "")[:120],
        due_date=due_date_dt,
        is_cram=is_cram,
        lead_time_days=lead_time_days,
        current_step=1,
    )
    track_doc = track.model_dump()
    track_doc["created_at"] = track_doc["created_at"].isoformat()
    track_doc["due_date"] = due_date_dt.isoformat() if due_date_dt else None
    track_doc["is_template"] = payload.is_template
    track_doc["step_completed_at"] = {}
    track_doc["current_bundle_index"] = 0

    now_analyze = datetime.now(timezone.utc)
    concepts_result = await extract_key_concepts(
        text,
        payload.track_type,
        gemini_generate=_gemini_generate,
        gemini_enabled=GEMINI_ENABLED,
    )
    if material_id:
        await db.materials.update_one(
            {"material_id": material_id},
            {
                "$set": {
                    "key_concepts": concepts_result.get("key_concepts") or [],
                    "language_split": concepts_result.get("language_split"),
                    "analyzed_at": now_analyze.isoformat(),
                }
            },
        )

    raw_units = concepts_to_knowledge_units(concepts_result.get("key_concepts") or [], payload.track_type)
    if not raw_units:
        raw_units = await _extract_knowledge_units(text, payload.track_type)
    now = datetime.now(timezone.utc)
    prepared = []
    for u in raw_units:
        unit = KnowledgeUnit(
            owner=unit_owner,
            track_id=track.track_id,
            material_id=material_id,
            track_type=payload.track_type,
            unit_type=u.get("unit_type", "word"),
            term=str(u.get("term", ""))[:300],
            meaning=u.get("meaning"),
            context=u.get("context"),
            key_concept_id=u.get("key_concept_id"),
            exact_term=u.get("exact_term"),
            simplified_explanation=u.get("simplified_explanation"),
            presentation=u.get("presentation"),
        )
        d = unit.model_dump()
        d["language"] = u.get("language") or detect_language(d.get("term", ""))
        d["created_at"] = d["created_at"].isoformat()
        d["next_due_at"] = now.isoformat()
        d["last_reviewed_at"] = None
        prepared.append(d)

    unit_docs = assign_bundle_indices(prepared)
    track_doc["bundle_count"] = bundle_count_for_units(unit_docs)
    await db.learning_tracks.insert_one(track_doc)
    if unit_docs:
        await db.knowledge_units.insert_many(unit_docs)

    return {
        "track": {**track.model_dump(), "bundle_count": track_doc["bundle_count"]},
        "unit_count": len(unit_docs),
        "bundle_count": track_doc["bundle_count"],
        "is_cram": is_cram,
        "urgency_tier": urgency_tier,
        "urgency_meta": TIER_META[urgency_tier],
        "lead_time_days": lead_time_days,
        "key_concepts": concepts_result.get("key_concepts") or [],
        "concept_count": concepts_result.get("concept_count", 0),
        "language_split": concepts_result.get("language_split"),
        "analysis_source": concepts_result.get("source"),
    }


async def _clone_track_for_student(source_track: dict, units: list, student_id: str, assigned_by: str) -> str:
    """Clone a template track + units for one student."""
    new_track_id = f"trk_{uuid.uuid4().hex[:10]}"
    new_track = {**source_track, "track_id": new_track_id}
    new_track["student_id"] = student_id
    new_track["owner"] = assigned_by or source_track.get("owner")
    new_track["assigned_by"] = assigned_by
    new_track["source_track_id"] = source_track.get("track_id")
    new_track["is_template"] = False
    new_track["current_step"] = 1
    new_track["current_bundle_index"] = 0
    new_track["step_completed_at"] = {}
    new_track["consecutive_fail_count"] = 0
    new_track["created_at"] = datetime.now(timezone.utc).isoformat()
    if isinstance(new_track.get("due_date"), datetime):
        new_track["due_date"] = new_track["due_date"].isoformat()
    await db.learning_tracks.insert_one(new_track)

    if units:
        new_units = []
        for u in units:
            nu = {**u, "unit_id": f"ku_{uuid.uuid4().hex[:10]}", "track_id": new_track_id, "owner": student_id}
            if isinstance(nu.get("created_at"), datetime):
                nu["created_at"] = nu["created_at"].isoformat()
            new_units.append(nu)
        await db.knowledge_units.insert_many(new_units)
    return new_track_id


@api_router.post("/tracks/{track_id}/assign")
async def assign_track_to_students(
    track_id: str,
    payload: TrackAssignRequest,
    user: Optional[User] = Depends(get_current_user),
):
    """Clone a track to one or more students, or all participants in a classroom room."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")

    track = await db.learning_tracks.find_one({"track_id": track_id}, {"_id": 0})
    if not track:
        raise HTTPException(status_code=404, detail="Track not found")

    is_teacher = (user.role or "") == "teacher" or track.get("owner") == user.user_id
    is_parent = track.get("assigned_by") == user.user_id or track.get("owner") == user.user_id
    if not is_teacher and not is_parent:
        links = await db.family_links.find(
            {"parent_user_id": user.user_id, "status": "active"}, {"_id": 0}
        ).to_list(20)
        if not links:
            raise HTTPException(status_code=403, detail="Not allowed to assign this track")

    student_ids: List[str] = list(payload.student_ids or [])

    if payload.room_code:
        code = payload.room_code.upper()
        room = await db.classrooms.find_one({"room_code": code, "teacher_id": user.user_id}, {"_id": 0})
        if not room:
            raise HTTPException(status_code=404, detail="Classroom not found")
        for p in room.get("participants", []):
            sid = p.get("owner")
            if sid and sid not in student_ids:
                student_ids.append(sid)

    if not student_ids:
        raise HTTPException(status_code=400, detail="No students selected — pick roster students or a room with participants")

    for sid in student_ids:
        if sid == user.user_id:
            continue
        link = await db.family_links.find_one(
            {"parent_user_id": user.user_id, "kid_owner_id": sid, "status": "active"}, {"_id": 0}
        )
        if not link and user.role != "teacher":
            raise HTTPException(status_code=403, detail=f"Not linked to student {sid}")

    units = await db.knowledge_units.find({"track_id": track_id}, {"_id": 0}).to_list(500)
    assigned = []
    for sid in student_ids:
        if sid == user.user_id:
            continue
        new_id = await _clone_track_for_student(track, units, sid, user.user_id)
        assigned.append({"student_id": sid, "track_id": new_id})

    if payload.room_code and assigned:
        await db.classrooms.update_one(
            {"room_code": payload.room_code.upper()},
            {"$set": {"assigned_track_id": track_id, "last_assignment": assigned}},
        )

    return {"assigned": assigned, "source_track_id": track_id, "count": len(assigned)}


async def _maybe_skip_locked_step6(track: dict) -> dict:
    """Step 6 is optional — unstick tracks waiting on the old 24h gate."""
    if track.get("track_type") != "reading_dictation":
        return track
    if int(track.get("current_step") or 1) != 6:
        return track
    completed = parse_step_completed_at(track.get("step_completed_at"))
    if 5 in completed and 6 not in completed:
        await db.learning_tracks.update_one(
            {"track_id": track["track_id"]},
            {"$set": {"current_step": 7}},
        )
        return {**track, "current_step": 7}
    return track


@api_router.get("/tracks/{track_id}/journey")
async def get_track_journey(track_id: str, owner: str = Depends(get_owner_id)):
    track = await db.learning_tracks.find_one({"track_id": track_id, **_track_access_filter(owner)}, {"_id": 0})
    if not track:
        raise HTTPException(status_code=404, detail="Track not found")
    track = await _maybe_skip_locked_step6(track)
    tt = track.get("track_type", "reading_dictation")
    if tt not in TRACK_STEP_TABLES:
        raise HTTPException(status_code=400, detail="Journey map not available for this track type")
    return journey_status(track)


@api_router.post("/tracks/{track_id}/step-battle")
async def start_step_battle(
    track_id: str,
    step: Optional[int] = None,
    owner: str = Depends(get_owner_id),
):
    """Start a journey step battle — reading or recitation track (G1–G20 + READ/HL)."""
    track = await db.learning_tracks.find_one({"track_id": track_id, **_track_access_filter(owner)}, {"_id": 0})
    if not track:
        raise HTTPException(status_code=404, detail="Track not found")
    track = await _maybe_skip_locked_step6(track)
    track_type = track.get("track_type", "reading_dictation")
    if track_type not in TRACK_STEP_TABLES:
        raise HTTPException(status_code=400, detail="Step battles require a journey track (reading/recital/quiz/exam)")

    max_step = max_step_for_track(track_type)
    play_step = int(step or track.get("current_step") or 1)
    play_step = max(1, min(play_step, max_step))

    lock = step_lock_status(track, play_step)
    if not lock.get("unlocked"):
        raise HTTPException(status_code=423, detail=lock.get("message") or "Step locked")

    units = await db.knowledge_units.find({"track_id": track_id}, {"_id": 0}).to_list(500)
    if not units:
        raise HTTPException(status_code=400, detail="No units in track")

    bundle_idx = int(track.get("current_bundle_index") or 0)
    bundle_units = units_for_bundle(units, bundle_idx)
    if not bundle_units:
        raise HTTPException(status_code=400, detail="No units in current bundle")

    performance = await _fetch_game_mode_performance(owner)

    title = track.get("scope_description") or {
        "recital_dictation": "Recitation",
        "quiz": "Quiz",
        "exam": "Exam",
    }.get(track_type, "Reading Dictation")
    if int(track.get("bundle_count") or 1) > 1:
        title += f" (Bundle {bundle_idx + 1}/{track.get('bundle_count')})"

    game = generate_step_game(bundle_units, play_step, title, performance=performance, track_type=track_type)
    unit_ids = [u["unit_id"] for u in bundle_units]

    return {
        "game": game,
        "unit_ids": unit_ids,
        "track_id": track_id,
        "journey_step": play_step,
        "bundle_index": bundle_idx,
        "journey": journey_status(track),
        "lock": lock,
    }


@api_router.get("/progress-snapshot")
async def get_progress_snapshot(owner: str = Depends(get_owner_id)):
    doc = await db.progress_snapshots.find_one({"owner": owner}, {"_id": 0})
    return doc or {"owner": owner, "snapshot": None}


@api_router.post("/progress-snapshot")
async def save_progress_snapshot(payload: ProgressSnapshotSave, owner: str = Depends(get_owner_id)):
    now = datetime.now(timezone.utc).isoformat()
    doc = {
        "owner": owner,
        "track_id": payload.track_id,
        "step_number": payload.step_number,
        "unit_ids": payload.unit_ids,
        "game": payload.game,
        "progress": payload.progress,
        "updated_at": now,
    }
    await db.progress_snapshots.update_one({"owner": owner}, {"$set": doc}, upsert=True)
    return {"ok": True, "updated_at": now}


@api_router.delete("/progress-snapshot")
async def clear_progress_snapshot(owner: str = Depends(get_owner_id)):
    await db.progress_snapshots.delete_one({"owner": owner})
    return {"ok": True}


async def _advance_journey_step(track_id: str, step: int, correct: int, wrong: int) -> dict:
    track = await db.learning_tracks.find_one({"track_id": track_id}, {"_id": 0})
    track_type = (track or {}).get("track_type", "reading_dictation")
    if not track or track_type not in TRACK_STEP_TABLES:
        return {}
    total = correct + wrong
    passed = score_meets_step_threshold(correct, total, step, track_type)
    now_iso = datetime.now(timezone.utc).isoformat()
    completed = dict(track.get("step_completed_at") or {})
    cfg = get_step_table(track_type).get(step, {})
    cur = int(track.get("current_step") or 1)
    max_step = max_step_for_track(track_type)
    is_replay = (
        _step_was_completed(completed, step)
        or step < cur
        or track.get("status") == "completed"
    )

    if is_replay:
        result = {
            "passed": passed,
            "step": step,
            "accuracy_pct": round((correct / total) * 100, 1) if total else 0,
            "pass_threshold_pct": cfg.get("pass_pct", 0),
            "replay": True,
            "current_step": cur,
            "status": track.get("status", "active"),
        }
        if passed:
            completed[str(step)] = now_iso
            await db.learning_tracks.update_one(
                {"track_id": track_id},
                {"$set": {"step_completed_at": completed}},
            )
        return result

    updates: dict = {}

    if passed:
        completed[str(step)] = now_iso
        if cfg.get("optional") and step < cur:
            next_s = cur
        else:
            next_s = next_step_after_pass(step, track_type)
        updates = {
            "step_completed_at": completed,
            "current_step": next_s,
            "consecutive_fail_count": 0,
        }
        if step >= max_step:
            bundle_idx = int(track.get("current_bundle_index") or 0)
            bundle_total = int(track.get("bundle_count") or 1)
            if bundle_idx + 1 < bundle_total:
                updates["current_bundle_index"] = bundle_idx + 1
                updates["current_step"] = 1
                updates["step_completed_at"] = {}
            else:
                updates["status"] = "completed"
                updates["current_step"] = max_step
        if track_type == "quiz" and step == 9:
            lock_h = random.randint(48, 72)
            slh = dict(track.get("step_lock_hours") or {})
            slh["10"] = lock_h
            updates["step_lock_hours"] = slh
    else:
        fail_count = int(track.get("consecutive_fail_count") or 0) + 1
        updates = {"consecutive_fail_count": fail_count}
        back = step_back_on_fail(step, correct, total, track_type)
        if back:
            updates["current_step"] = back
        elif track_type == "reading_dictation" and step == 6:
            updates["current_step"] = 4
        elif fail_count >= 2 and step > 1 and not (
            track_type in ("reading_dictation", "recital_dictation") and step == 10
        ):
            updates["current_step"] = max(1, step - 2)

    await db.learning_tracks.update_one({"track_id": track_id}, {"$set": updates})
    result = {
        "passed": passed,
        "step": step,
        "accuracy_pct": round((correct / total) * 100, 1) if total else 0,
        "pass_threshold_pct": cfg.get("pass_pct", 0),
        **updates,
    }
    if passed and updates.get("status") == "completed":
        result["graduated"] = True
        result["boss_followup_days"] = GRADUATION_BOSS_DELAY_DAYS
        await _schedule_graduation_boss_followup({**track, **updates}, datetime.now(timezone.utc))
    return result


@api_router.get("/tracks")
async def list_tracks(
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
    kid_owner_id: Optional[str] = None,
):
    if kid_owner_id and user:
        link = await db.family_links.find_one(
            {"parent_user_id": user.user_id, "kid_owner_id": kid_owner_id, "status": "active"},
            {"_id": 0},
        )
        if not link:
            raise HTTPException(status_code=403, detail="Not linked to this child")
        query = _student_tracks_filter(kid_owner_id)
    elif user and owner == user.user_id:
        kid_ids = await _parent_linked_kid_ids(user.user_id)
        if not kid_ids:
            query = {"track_id": {"$in": []}}
        elif kid_owner_id:
            if kid_owner_id not in kid_ids:
                raise HTTPException(status_code=403, detail="Not linked to this child")
            query = _student_tracks_filter(kid_owner_id)
        else:
            if len(kid_ids) == 1:
                query = _student_tracks_filter(kid_ids[0])
            else:
                raise HTTPException(status_code=400, detail="kid_owner_id required when multiple children linked")
    else:
        query = _student_tracks_filter(owner)

    # Kids should not see teacher/parent template tracks — only their playable copies.
    if not user or (owner.startswith("guest_") and owner != (user.user_id if user else "")):
        query = {"$and": [query, {"$or": [{"is_template": {"$ne": True}}, {"is_template": {"$exists": False}}]}]}

    tracks = await db.learning_tracks.find(query, {"_id": 0}).sort("created_at", -1).to_list(100)
    now = datetime.now(timezone.utc)
    out = []
    for t in tracks:
        units = await db.knowledge_units.find({"track_id": t["track_id"]}, {"_id": 0}).to_list(500)
        strengths = [_unit_memory_strength(u, now) for u in units]
        readiness = round(sum(strengths) / len(strengths), 1) if strengths else 0.0
        weak_sorted = sorted(units, key=lambda u: _unit_memory_strength(u, now))
        weak = [u for u in weak_sorted if _unit_memory_strength(u, now) < MEMORY_THRESHOLD][:5]
        t["readiness_percent"] = readiness
        t["unit_count"] = len(units)
        t["weak_units"] = [
            {"unit_id": u["unit_id"], "term": u["term"], "memory_strength": _unit_memory_strength(u, now)}
            for u in weak
        ]
        _attach_tier(t, now)
        out.append(t)
    return out


def _bucket_units_by_due_date(units: list, days_ahead: int) -> dict:
    """Group knowledge units by the calendar date their next review is due —
    powers the "daily review load" forecast on the calendar (Section: parent/
    teacher calendar request)."""
    now = datetime.now(timezone.utc)
    buckets: dict = {}
    for u in units:
        due_dt = _parse_dt(u.get("next_due_at"))
        if not due_dt:
            continue
        delta_days = (due_dt.date() - now.date()).days
        if delta_days < 0 or delta_days > days_ahead:
            continue
        key = due_dt.date().isoformat()
        buckets[key] = buckets.get(key, 0) + 1
    return buckets


@api_router.get("/calendar")
async def get_calendar(
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
    kid_owner_id: Optional[str] = None,
    days_ahead: int = 21,
):
    """Deadlines (quiz/exam due dates) + daily review-load forecast for the
    Parent (and self-practice Kid) dashboard calendar. Mirrors the same
    access-control pattern as GET /tracks."""
    if kid_owner_id and user:
        link = await db.family_links.find_one(
            {"parent_user_id": user.user_id, "kid_owner_id": kid_owner_id, "status": "active"},
            {"_id": 0},
        )
        if not link:
            raise HTTPException(status_code=403, detail="Not linked to this child")
        query = _student_tracks_filter(kid_owner_id)
    elif user and owner == user.user_id:
        kid_ids = await _parent_linked_kid_ids(user.user_id)
        if not kid_ids:
            query = {"track_id": {"$in": []}}
        elif kid_owner_id:
            if kid_owner_id not in kid_ids:
                raise HTTPException(status_code=403, detail="Not linked to this child")
            query = _student_tracks_filter(kid_owner_id)
        else:
            if len(kid_ids) == 1:
                query = _student_tracks_filter(kid_ids[0])
            else:
                raise HTTPException(status_code=400, detail="kid_owner_id required when multiple children linked")
    else:
        query = _student_tracks_filter(owner)

    tracks = await db.learning_tracks.find(query, {"_id": 0}).to_list(200)
    now = datetime.now(timezone.utc)
    horizon = now + timedelta(days=days_ahead)

    events = []
    tids = [t["track_id"] for t in tracks]
    for t in tracks:
        due_dt = _parse_dt(t.get("due_date")) if t.get("due_date") else None
        if due_dt and now <= due_dt <= horizon:
            events.append({
                "date": due_dt.date().isoformat(),
                "type": "deadline",
                "track_id": t["track_id"],
                "title": t.get("title", "Untitled"),
                "track_type": t.get("track_type"),
                "urgency_tier": compute_urgency_tier(t.get("track_type"), due_dt, now),
            })

    units = await db.knowledge_units.find({"track_id": {"$in": tids}}, {"_id": 0}).to_list(2000) if tids else []
    daily_load = _bucket_units_by_due_date(units, days_ahead)

    return {
        "events": sorted(events, key=lambda e: e["date"]),
        "daily_load": [{"date": d, "units_due": c} for d, c in sorted(daily_load.items())],
        "generated_at": now.isoformat(),
    }


@api_router.get("/tracks/{track_id}")
async def get_track(track_id: str, owner: str = Depends(get_owner_id)):
    t = await db.learning_tracks.find_one({"track_id": track_id, **_track_access_filter(owner)}, {"_id": 0})
    if not t:
        raise HTTPException(status_code=404, detail="Track not found")
    units = await db.knowledge_units.find({"track_id": track_id}, {"_id": 0}).to_list(500)
    now = datetime.now(timezone.utc)
    for u in units:
        u["memory_strength"] = _unit_memory_strength(u, now)
    strengths = [u["memory_strength"] for u in units]
    t["readiness_percent"] = round(sum(strengths) / len(strengths), 1) if strengths else 0.0
    t["units"] = units
    _attach_tier(t, now)
    return t


@api_router.delete("/tracks/{track_id}")
async def delete_track(track_id: str, owner: str = Depends(get_owner_id)):
    """Remove a learning track and its knowledge units."""
    track = await db.learning_tracks.find_one({"track_id": track_id, **_track_access_filter(owner)}, {"_id": 0})
    if not track:
        raise HTTPException(status_code=404, detail="Track not found")
    await db.knowledge_units.delete_many({"track_id": track_id, "owner": owner})
    await db.learning_tracks.delete_one({"track_id": track_id, "owner": owner})
    await db.progress_snapshots.delete_one({"owner": owner, "track_id": track_id})
    return {"ok": True, "track_id": track_id}


@api_router.post("/reviews")
async def record_review(payload: ReviewCreate, owner: str = Depends(get_owner_id)):
    await assert_can_collect_learning_data(db, owner)
    unit = await db.knowledge_units.find_one({"unit_id": payload.unit_id, "owner": owner}, {"_id": 0})
    if not unit:
        raise HTTPException(status_code=404, detail="Knowledge unit not found")

    quality = payload.quality_override if payload.quality_override is not None else \
        quality_from_response(payload.correct, payload.response_time_ms)

    new_ease, new_interval, new_reps = sm2_update(
        unit.get("ease_factor", DEFAULT_EASE_FACTOR),
        unit.get("interval_days", 0) or 0,
        unit.get("repetitions", 0),
        quality,
    )
    now = datetime.now(timezone.utc)
    next_due = now + timedelta(days=new_interval)

    update = {
        "ease_factor": new_ease,
        "interval_days": new_interval,
        "repetitions": new_reps,
        "review_count": unit.get("review_count", 0) + 1,
        "last_quality": quality,
        "last_reviewed_at": now.isoformat(),
        "next_due_at": next_due.isoformat(),
    }
    await db.knowledge_units.update_one({"unit_id": payload.unit_id}, {"$set": update})

    await db.reviews.insert_one({
        "review_id": f"rv_{uuid.uuid4().hex[:10]}",
        "unit_id": payload.unit_id,
        "owner": owner,
        "correct": payload.correct,
        "response_time_ms": payload.response_time_ms,
        "quality": quality,
        "reviewed_at": now.isoformat(),
    })

    await _maybe_resolve_follow_ups(owner, payload.unit_id, now)

    updated_unit = {**unit, **update}
    updated_unit["memory_strength"] = 100.0  # freshly reviewed, decay starts now
    return updated_unit


@api_router.get("/daily-queue")
async def daily_queue(owner: str = Depends(get_owner_id), max_units: int = 12):
    tracks = await db.learning_tracks.find(_student_tracks_filter(owner), {"_id": 0}).to_list(200)
    if not tracks:
        return {"due": [], "new": [], "total_due_count": 0, "total_new_count": 0, "total_duration_estimate_min": 0}

    now = datetime.now(timezone.utc)
    tier_rank = {"survival": 0, "emergency": 1, "cram": 2, "standard": 3}
    track_map = {}
    for t in tracks:
        tier = compute_urgency_tier(t.get("track_type"), t.get("due_date"), now)
        track_map[t["track_id"]] = {**t, "urgency_tier": tier}

    units = await db.knowledge_units.find({"track_id": {"$in": list(track_map.keys())}}, {"_id": 0}).to_list(2000)

    followup_due_ids = set()
    followups = await db.follow_up_tasks.find({"owner": owner, "status": "pending"}, {"_id": 0}).to_list(200)
    for f in followups:
        sched = _parse_dt(f.get("scheduled_date"))
        if sched and sched <= now:
            followup_due_ids.update(f.get("unit_ids", []))

    due, new = [], []
    for u in units:
        strength = _unit_memory_strength(u, now)
        u["memory_strength"] = strength
        u["urgency_tier"] = track_map.get(u["track_id"], {}).get("urgency_tier", "standard")
        u["is_cram"] = u["urgency_tier"] != "standard"  # kept for backward-compat consumers
        u["is_followup"] = u["unit_id"] in followup_due_ids
        if u.get("review_count", 0) == 0:
            new.append(u)
        elif strength < MEMORY_THRESHOLD or u["is_followup"]:
            due.append(u)

    # Priority per Section 6.4/12.2, with the Section 11.3 rule on top: a due
    # Consolidation Follow-up always jumps the queue, ahead of urgency tier, because
    # cram/emergency memory decays fastest right after the deadline.
    # self_practice tracks sort after deadline-assigned work (Section 12.2).
    def _queue_key(u):
        tr = track_map.get(u["track_id"], {})
        sp = 1 if tr.get("track_type") == "self_practice" else 0
        return (
            0 if u["is_followup"] else 1,
            sp,
            tier_rank.get(u["urgency_tier"], 3),
            u["memory_strength"],
        )

    due.sort(key=_queue_key)
    new.sort(key=lambda u: _queue_key(u)[1:])

    selected_due = due[:max_units]
    remaining = max(0, max_units - len(selected_due))
    selected_new = new[:remaining]
    est_minutes = round((len(selected_due) + len(selected_new)) * 0.75, 1)

    # Per-track readiness rings for Kid Dashboard (Section 7.2)
    track_rings = []
    for t in tracks:
        t_units = [u for u in units if u["track_id"] == t["track_id"]]
        strengths = [_unit_memory_strength(u, now) for u in t_units]
        readiness = round(sum(strengths) / len(strengths), 1) if strengths else 0.0
        due_count = sum(1 for u in t_units if u.get("review_count", 0) > 0 and _unit_memory_strength(u, now) < MEMORY_THRESHOLD)
        track_rings.append({
            "track_id": t["track_id"],
            "track_type": t.get("track_type"),
            "scope_description": t.get("scope_description", ""),
            "readiness_percent": readiness,
            "units_due": due_count,
            "unit_count": len(t_units),
            "urgency_tier": track_map.get(t["track_id"], {}).get("urgency_tier", "standard"),
        })

    bundle_summary = []
    for t in tracks:
        tid = t["track_id"]
        t_units = [u for u in units if u["track_id"] == tid]
        bundle_total = int(t.get("bundle_count") or 1)
        current_bi = int(t.get("current_bundle_index") or 0)
        for bi in range(bundle_total):
            bu = [u for u in t_units if u.get("bundle_index", 0) == bi]
            if not bu and bi > 0:
                continue
            strengths = [_unit_memory_strength(u, now) for u in bu]
            due_n = sum(
                1 for u in bu
                if u.get("review_count", 0) > 0 and _unit_memory_strength(u, now) < MEMORY_THRESHOLD
            )
            bundle_summary.append({
                "track_id": tid,
                "track_type": t.get("track_type"),
                "scope_description": t.get("scope_description", ""),
                "bundle_index": bi,
                "bundle_total": bundle_total,
                "unit_count": len(bu),
                "units_due": due_n,
                "readiness_percent": round(sum(strengths) / len(strengths), 1) if strengths else 0.0,
                "is_current": bi == current_bi,
            })

    return {
        "due": selected_due,
        "new": selected_new,
        "total_due_count": len(due),
        "total_new_count": len(new),
        "total_duration_estimate_min": est_minutes,
        "track_rings": track_rings,
        "bundle_summary": bundle_summary,
    }


@api_router.get("/follow-ups")
async def list_follow_ups(
    owner: str = Depends(get_owner_id),
    user: Optional[User] = Depends(get_current_user),
    kid_owner_id: Optional[str] = None,
):
    """Section 10.4/11.5: 1-2 day post-cram consolidation retests. 'due' = ready to
    retest now, 'upcoming' = scheduled but not yet due, 'resolved' = done or escalated."""
    data_owner = await _resolve_child_data_owner(user, owner, kid_owner_id)
    if data_owner is None:
        return {"due": [], "upcoming": [], "resolved": []}
    now = datetime.now(timezone.utc)
    docs = await db.follow_up_tasks.find({"owner": data_owner}, {"_id": 0}).sort("scheduled_date", 1).to_list(200)
    due, upcoming, resolved = [], [], []
    for d in docs:
        if d.get("status") != "pending":
            resolved.append(d)
            continue
        sched = _parse_dt(d.get("scheduled_date"))
        (due if (sched and sched <= now) else upcoming).append(d)
    return {"due": due, "upcoming": upcoming, "resolved": resolved[:20]}


# ============== CONSOLIDATION ENERGY (Brain HP) ==============
# Section 10.1 forgetting-curve framing: HP is not a battle resource that gates
# you — it's a visualisation of "brain consolidation" (Rest & Sleep). Depletes
# when you battle; refills naturally over hours ~ mimicking real memory
# consolidation. Never blocks Core Loop (due units always playable); only paces
# Side Loop / free-play so kids intuit "your brain needs rest to lock things in".

MAX_ENERGY = 100.0
ENERGY_REFILL_PER_HOUR = 12.5  # → full refill in 8h (approximates one sleep cycle)
ENERGY_BASE_COST = 5.0
ENERGY_COST_PER_CHALLENGE = 1.5


async def _compute_energy(owner: str, now: datetime) -> dict:
    """Derived from game_sessions history over last 24h. No mutation needed."""
    since = now - timedelta(hours=24)
    sessions = await db.game_sessions.find(
        {"owner": owner, "completed_at": {"$gte": since.isoformat()}},
        {"_id": 0}
    ).sort("completed_at", 1).to_list(200)

    energy = MAX_ENERGY
    last_time = since
    for s in sessions:
        ts = _parse_dt(s.get("completed_at"))
        if not ts:
            continue
        # refill in the gap since previous event
        gap_h = max(0.0, (ts - last_time).total_seconds() / 3600.0)
        energy = min(MAX_ENERGY, energy + gap_h * ENERGY_REFILL_PER_HOUR)
        cost = ENERGY_BASE_COST + (s.get("correct", 0) + s.get("wrong", 0)) * ENERGY_COST_PER_CHALLENGE
        # Free-play battles cost slightly less (marked in mode)
        if s.get("mode") == "free_play":
            cost *= 0.6
        energy = max(0.0, energy - cost)
        last_time = ts

    # refill from last session up to now
    gap_h = max(0.0, (now - last_time).total_seconds() / 3600.0)
    energy = min(MAX_ENERGY, energy + gap_h * ENERGY_REFILL_PER_HOUR)

    full_at = None
    if energy < MAX_ENERGY:
        hours_needed = (MAX_ENERGY - energy) / ENERGY_REFILL_PER_HOUR
        full_at = (now + timedelta(hours=hours_needed)).isoformat()

    return {
        "energy": round(energy, 1),
        "max_energy": MAX_ENERGY,
        "refill_per_hour": ENERGY_REFILL_PER_HOUR,
        "full_refill_at": full_at,
        "pct": round(energy / MAX_ENERGY * 100, 1),
    }


@api_router.get("/energy")
async def get_energy(owner: str = Depends(get_owner_id)):
    return await _compute_energy(owner, datetime.now(timezone.utc))


# ============== HOME STATUS (Kid Dashboard aggregate) ==============
REGION_UNLOCK_RULES = [
    ("forest", 0, 0),
    ("plain", 1, 0),
    ("abyss", 3, 5),
    ("mount", 5, 10),
    ("tower", 8, 20),
]


@api_router.get("/home-status")
async def get_home_status(owner: str = Depends(get_owner_id)):
    now = datetime.now(timezone.utc)

    energy = await _compute_energy(owner, now)

    tracks = await db.learning_tracks.find(
        _student_tracks_filter(owner), {"_id": 0}
    ).to_list(200)

    tier_rank = {"survival": 0, "emergency": 1, "cram": 2, "standard": 3}
    boss_units_due = 0
    urgency_tier = "standard"
    track_rings = []

    if tracks:
        track_tier_map = {
            t["track_id"]: compute_urgency_tier(t.get("track_type"), t.get("due_date"), now)
            for t in tracks
        }
        units = await db.knowledge_units.find(
            {"track_id": {"$in": list(track_tier_map.keys())}}, {"_id": 0}
        ).to_list(2000)
        followup_due_ids = set()
        followups = await db.follow_up_tasks.find(
            {"owner": owner, "status": "pending"}, {"_id": 0}
        ).to_list(200)
        for f in followups:
            sched = _parse_dt(f.get("scheduled_date"))
            if sched and sched <= now:
                followup_due_ids.update(f.get("unit_ids", []))

        due_units = []
        new_units = []
        for u in units:
            ms = _unit_memory_strength(u, now)
            is_fu = u["unit_id"] in followup_due_ids
            if (u.get("review_count", 0) > 0 and ms < MEMORY_THRESHOLD) or is_fu:
                due_units.append(u)
            elif u.get("review_count", 0) == 0:
                new_units.append(u)

        boss_units_due = len(due_units) + min(len(new_units), 8)
        if boss_units_due > 0:
            relevant = due_units + new_units[:8]
            urgency_tier = min(
                (track_tier_map.get(u["track_id"], "standard") for u in relevant),
                key=lambda t: tier_rank.get(t, 3),
            )

        for t in tracks:
            t_units = [u for u in units if u["track_id"] == t["track_id"]]
            strengths = [_unit_memory_strength(u, now) for u in t_units]
            readiness = round(sum(strengths) / len(strengths), 1) if strengths else 0.0
            due_count = sum(
                1 for u in t_units
                if u.get("review_count", 0) > 0 and _unit_memory_strength(u, now) < MEMORY_THRESHOLD
            )
            track_rings.append({
                "track_id": t["track_id"],
                "track_type": t.get("track_type"),
                "scope_description": t.get("scope_description", ""),
                "readiness_percent": readiness,
                "units_due": due_count,
                "unit_count": len(t_units),
                "urgency_tier": track_tier_map.get(t["track_id"], "standard"),
            })

    family_link_count = await _active_family_link_count(owner)

    total_sessions = await db.game_sessions.count_documents({"owner": owner})
    total_tracks = len(tracks)
    boss_defeats = await db.game_sessions.count_documents(
        {"owner": owner, "correct": {"$gte": 5}}
    )

    regions_unlocked = {
        rid: (total_tracks >= min_tracks or (min_defeats > 0 and boss_defeats >= min_defeats))
        for rid, min_tracks, min_defeats in REGION_UNLOCK_RULES
    }

    return {
        "energy": energy,
        "boss_status": {
            "ready": boss_units_due > 0,
            "units_due": boss_units_due,
            "urgency_tier": urgency_tier,
            "urgency_meta": TIER_META.get(urgency_tier, TIER_META["standard"]),
        },
        "track_rings": track_rings,
        "family_link_count": family_link_count,
        "regions_unlocked": regions_unlocked,
        "stats": {
            "total_sessions": total_sessions,
            "total_tracks": total_tracks,
            "boss_defeats": boss_defeats,
        },
        "streak": await db.streaks.find_one({"owner": owner}, {"_id": 0}) or {
            "current_streak": 0, "longest_streak": 0,
        },
        "data_consent": await get_data_consent_status(db, owner),
    }


# ============== FREE PLAY (Side Loop) ==============
# Section 12: pure practice/fun mode. Doesn't touch knowledge_units, doesn't
# affect memory_strength, doesn't reschedule follow-ups, doesn't feed readiness.
# Available whenever: on rest days (no due content), or when the kid just feels
# like playing with a favourite past worksheet.

class FreePlayCreate(BaseModel):
    material_id: Optional[str] = None
    text: Optional[str] = None
    mode: Optional[str] = "quiz"


@api_router.post("/free-play")
async def start_free_play(payload: FreePlayCreate, owner: str = Depends(get_owner_id)):
    text = None
    if payload.material_id:
        mat = await db.materials.find_one({"material_id": payload.material_id}, {"_id": 0})
        if not mat:
            raise HTTPException(status_code=404, detail="Material not found")
        if mat.get("owner") not in (owner, "guest"):
            raise HTTPException(status_code=403, detail="Not your material")
        text = mat.get("text", "")
    elif payload.text:
        text = payload.text
    else:
        raise HTTPException(status_code=400, detail="Provide material_id or text")

    if not text or len(text.strip()) < 8:
        raise HTTPException(status_code=400, detail="Material too short for a game")

    req = GenerateGameRequest(
        text=text, mode=payload.mode or "quiz", subject="general", difficulty=2
    )
    result = await generate_game(req)
    return {**result, "is_free_play": True}




# ============== FAMILY LINK / CONSENT / PRIVACY ENDPOINTS ==============
@api_router.get("/kid-device/family-code")
async def get_kid_family_code(
    x_guest_id: Optional[str] = Header(default=None, alias="X-Guest-Id"),
):
    if not x_guest_id:
        raise HTTPException(status_code=400, detail="Device id required (X-Guest-Id)")
    code = await _ensure_family_code(x_guest_id)
    kid_owner = kid_owner_from_device(x_guest_id)
    return {"family_code": code, "kid_device_id": x_guest_id, "kid_owner_id": kid_owner}


@api_router.post("/family-links/by-code")
async def link_family_by_code(
    payload: FamilyLinkByCodeCreate,
    user: Optional[User] = Depends(get_current_user),
):
    """Parent-initiated in-person link — immediate active (parent already at kid device)."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    fc = await _resolve_family_code(payload.family_code.strip())
    kid_device_id = fc["kid_device_id"]
    kid_owner_id = kid_owner_from_device(kid_device_id)

    existing = await db.family_links.find_one(
        {"parent_user_id": user.user_id, "kid_owner_id": kid_owner_id, "status": "active"},
        {"_id": 0},
    )
    if existing:
        return existing

    link = FamilyLink(
        parent_user_id=user.user_id,
        kid_owner_id=kid_owner_id,
        kid_device_id=kid_device_id,
        permission_level=payload.permission_level or "manage_tracks",
        status="active",
    )
    doc = link.model_dump()
    doc["created_at"] = doc["created_at"].isoformat()
    await db.family_links.insert_one(doc)

    now = datetime.now(timezone.utc)
    consent = ConsentRecord(
        parent_user_id=user.user_id,
        kid_owner_id=kid_owner_id,
        consent_type="data_collection",
        granted=True,
        method="in_person",
        granted_at=now,
        link_id=link.link_id,
    )
    cdoc = consent.model_dump()
    cdoc["created_at"] = cdoc["created_at"].isoformat()
    cdoc["granted_at"] = now.isoformat()
    await db.consent_records.insert_one(cdoc)
    return link


@api_router.get("/parent/child-device")
async def parent_child_device(
    kid_owner_id: str,
    user: User = Depends(get_current_user),
):
    if not user or user.role != "parent":
        raise HTTPException(status_code=403, detail="Parent account required")
    await _require_parent_child_link(user, kid_owner_id)
    return await get_device_profile(db, kid_owner_id)


@api_router.post("/parent/child-subscription")
async def parent_update_child_subscription(
    payload: ChildSubscriptionUpdate,
    user: User = Depends(get_current_user),
):
    """Parent enables premium on a linked child's device (IAP wiring later — early access)."""
    if not user or user.role != "parent":
        raise HTTPException(status_code=403, detail="Parent account required")
    await _require_parent_child_link(user, payload.kid_owner_id)
    tier = payload.subscription_tier
    if tier not in ("free", "premium"):
        raise HTTPException(status_code=400, detail="tier must be free or premium")
    now = datetime.now(timezone.utc).isoformat()
    await db.device_profiles.update_one(
        {"owner": payload.kid_owner_id},
        {"$set": {"subscription_tier": tier, "updated_at": now}, "$setOnInsert": {"owner": payload.kid_owner_id, "coins": 0}},
        upsert=True,
    )
    return {"kid_owner_id": payload.kid_owner_id, "subscription_tier": tier}


@api_router.post("/family-links/invite")
async def invite_family_link(
    payload: FamilyLinkInviteCreate,
    x_guest_id: Optional[str] = Header(default=None, alias="X-Guest-Id"),
):
    """Kid-initiated remote invite — pending until parent confirms via email token."""
    device_id = payload.kid_device_id or x_guest_id
    if not device_id:
        raise HTTPException(status_code=400, detail="kid_device_id required")
    kid_owner_id = kid_owner_from_device(device_id)
    email = payload.parent_email.strip().lower()
    if not email or "@" not in email:
        raise HTTPException(status_code=400, detail="Valid parent email required")

    pending = await db.family_links.find_one(
        {"kid_owner_id": kid_owner_id, "parent_email": email, "status": "pending"},
        {"_id": 0},
    )
    if pending:
        return {"ok": True, "status": "pending", "message": "Invite already sent"}

    token = uuid.uuid4().hex
    link = FamilyLink(
        parent_user_id="",
        kid_owner_id=kid_owner_id,
        kid_device_id=device_id.removeprefix("guest_"),
        permission_level="manage_tracks",
        status="pending",
        parent_email=email,
    )
    ldoc = link.model_dump()
    ldoc["created_at"] = ldoc["created_at"].isoformat()
    await db.family_links.insert_one(ldoc)

    record = ConsentRecord(
        parent_user_id="",
        kid_owner_id=kid_owner_id,
        parent_email=email,
        link_id=link.link_id,
        consent_type="data_collection",
        granted=False,
        method="email_link",
        consent_token=token,
    )
    rdoc = record.model_dump()
    rdoc["created_at"] = rdoc["created_at"].isoformat()
    await db.consent_records.insert_one(rdoc)

    confirm_url = f"{os.environ.get('APP_BASE_URL', 'http://localhost:3000')}/consent/confirm?token={token}"
    await _send_consent_email(email, confirm_url, "data_collection")
    return {"ok": True, "status": "pending", "message": "Invite sent — waiting for parent email confirmation"}


@api_router.post("/family-links")
async def create_family_link(
    payload: FamilyLinkCreate,
    user: Optional[User] = Depends(get_current_user),
):
    """Legacy direct link — prefer /family-links/by-code for new integrations."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    link = FamilyLink(
        parent_user_id=user.user_id,
        kid_owner_id=payload.kid_owner_id,
        permission_level=payload.permission_level or "manage_tracks",
        status="active",
    )
    doc = link.model_dump()
    doc["created_at"] = doc["created_at"].isoformat()
    await db.family_links.insert_one(doc)
    return link


@api_router.get("/family-links")
async def list_family_links(user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    links = await db.family_links.find(
        {"parent_user_id": user.user_id, "status": "active"}, {"_id": 0}
    ).to_list(50)
    out = []
    for l in links:
        kid = l["kid_owner_id"]
        tracks = await db.learning_tracks.find(_student_tracks_filter(kid), {"_id": 0}).to_list(50)
        sessions = await db.game_sessions.find({"owner": kid}, {"_id": 0}).sort("completed_at", -1).to_list(5)
        out.append({**l, "track_count": len(tracks), "recent_sessions": sessions})
    return out


@api_router.get("/family-links/orphaned")
async def list_orphaned_devices(user: Optional[User] = Depends(get_current_user)):
    """Tracks from previously linked devices no longer active (device change edge case)."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    active_kids = {
        l["kid_owner_id"]
        for l in await db.family_links.find(
            {"parent_user_id": user.user_id, "status": "active"}, {"kid_owner_id": 1, "_id": 0}
        ).to_list(100)
    }
    historical = await db.family_links.find(
        {"parent_user_id": user.user_id}, {"_id": 0}
    ).to_list(200)
    orphaned = []
    seen = set()
    for l in historical:
        kid = l.get("kid_owner_id")
        if kid and kid not in active_kids and kid not in seen:
            seen.add(kid)
            count = await db.learning_tracks.count_documents({"student_id": kid})
            if count:
                orphaned.append({"kid_owner_id": kid, "track_count": count, "last_linked": l.get("created_at")})
    return orphaned


@api_router.post("/consent/request")
async def request_consent(
    payload: ConsentRequest,
    user: Optional[User] = Depends(get_current_user),
):
    """Step 1 of verifiable parental consent: parent (logged in) requests consent
    for a specific kid + a specific consent_type. We email them a one-click
    confirm link rather than just flipping a flag — a click from inside the
    already-logged-in session doesn't prove the *parent* saw and agreed to the
    current policy text, a separate emailed confirmation does."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    await _require_parent_child_link(user, payload.kid_owner_id)
    token = uuid.uuid4().hex
    record = ConsentRecord(
        parent_user_id=user.user_id,
        kid_owner_id=payload.kid_owner_id,
        consent_type=payload.consent_type,
        granted=False,
        method="email_link",
        consent_token=token,
    )
    doc = record.model_dump()
    doc["created_at"] = doc["created_at"].isoformat()
    await db.consent_records.insert_one(doc)
    confirm_url = f"{os.environ.get('APP_BASE_URL', 'https://your-app.example.com')}/consent/confirm?token={token}"
    await _send_consent_email(payload.parent_email, confirm_url, payload.consent_type)
    return {"consent_id": record.consent_id, "status": "pending", "confirm_url": confirm_url}


@api_router.get("/consent/confirm")
async def confirm_consent(token: str, user: Optional[User] = Depends(get_current_user)):
    """Public endpoint — parent confirms via emailed link. Activates pending family_link."""
    rec = await db.consent_records.find_one({"consent_token": token, "granted": False}, {"_id": 0})
    if not rec:
        raise HTTPException(status_code=404, detail="Invalid or already-used consent link")
    now = datetime.now(timezone.utc)

    parent_user_id = rec.get("parent_user_id") or (user.user_id if user else "")
    if not parent_user_id and rec.get("parent_email"):
        existing = await db.users.find_one({"email": rec["parent_email"]}, {"_id": 0})
        if existing:
            parent_user_id = existing["user_id"]
        elif user:
            parent_user_id = user.user_id
        else:
            parent_user_id = f"user_{uuid.uuid4().hex[:12]}"
            await db.users.insert_one({
                "user_id": parent_user_id,
                "email": rec["parent_email"],
                "name": rec["parent_email"].split("@")[0],
                "picture": None,
                "role": "parent",
                "created_at": now.isoformat(),
            })

    await db.consent_records.update_one(
        {"consent_id": rec["consent_id"]},
        {"$set": {
            "granted": True,
            "granted_at": now.isoformat(),
            "consent_token": None,
            "parent_user_id": parent_user_id,
        }},
    )

    if rec.get("link_id"):
        await db.family_links.update_one(
            {"link_id": rec["link_id"], "status": "pending"},
            {"$set": {
                "status": "active",
                "parent_user_id": parent_user_id,
                "parent_email": rec.get("parent_email"),
            }},
        )
    elif rec.get("kid_owner_id"):
        await db.family_links.update_one(
            {"kid_owner_id": rec["kid_owner_id"], "status": "pending"},
            {"$set": {"status": "active", "parent_user_id": parent_user_id}},
        )

    return {
        "ok": True,
        "consent_type": rec["consent_type"],
        "kid_owner_id": rec["kid_owner_id"],
        "message": "Consent confirmed — family link is now active",
    }


@api_router.post("/consent/{consent_id}/revoke")
async def revoke_consent(consent_id: str, user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    rec = await db.consent_records.find_one({"consent_id": consent_id, "parent_user_id": user.user_id}, {"_id": 0})
    if not rec:
        raise HTTPException(status_code=404, detail="Consent record not found")
    await db.consent_records.update_one(
        {"consent_id": consent_id},
        {"$set": {"granted": False, "revoked_at": datetime.now(timezone.utc).isoformat()}},
    )
    return {"ok": True}


@api_router.get("/consent/status")
async def consent_status(kid_owner_id: str, user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    docs = await db.consent_records.find(
        {"kid_owner_id": kid_owner_id, "parent_user_id": user.user_id}, {"_id": 0}
    ).sort("created_at", -1).to_list(50)
    # latest record per consent_type wins
    latest = {}
    for d in docs:
        if d["consent_type"] not in latest:
            latest[d["consent_type"]] = d
    return latest


@api_router.delete("/account")
async def erase_account(owner: str = Depends(get_owner_id)):
    """GDPR-K / PDPO 'right to erasure' — hard delete, not a soft flag. Cascades
    through every collection keyed by owner. This is deliberately blunt; add a
    confirmation step in the UI before wiring a button to this."""
    collections = [
        "materials", "learning_tracks", "knowledge_units", "game_sessions",
        "follow_up_tasks", "consent_records", "family_links", "streaks", "email_preferences",
    ]
    deleted = {}
    for coll in collections:
        if coll == "consent_records":
            result = await db.consent_records.delete_many({"$or": [{"parent_user_id": owner}, {"kid_owner_id": owner}]})
        elif coll == "family_links":
            result = await db.family_links.delete_many({"$or": [{"parent_user_id": owner}, {"kid_owner_id": owner}]})
        elif coll == "email_preferences":
            result = await db.email_preferences.delete_many({"user_id": owner})
        else:
            result = await db[coll].delete_many({"owner": owner})
        deleted[coll] = result.deleted_count
    bs = await db.behavioral_signals.delete_many({"student_id": owner})
    deleted["behavioral_signals"] = bs.deleted_count
    return {"ok": True, "deleted": deleted}


@api_router.get("/account/export")
async def export_account(owner: str = Depends(get_owner_id)):
    """GDPR/PDPO 'right to access' — full dump of everything held about one owner."""
    collections = ["materials", "learning_tracks", "knowledge_units", "game_sessions", "follow_up_tasks"]
    export = {"owner": owner, "exported_at": datetime.now(timezone.utc).isoformat()}
    for coll in collections:
        export[coll] = await db[coll].find({"owner": owner}, {"_id": 0}).to_list(5000)
    return export


# ============== CLASSROOM (Teacher) ==============
class Classroom(BaseModel):
    room_code: str
    teacher_id: str
    name: str
    mode: str
    material_id: Optional[str] = None
    status: str = "waiting"  # waiting | active | finished
    participants: List[dict] = Field(default_factory=list)
    active_game: Optional[dict] = None
    expires_at: Optional[str] = None
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


class ClassroomCreate(BaseModel):
    name: str
    mode: str
    material_id: Optional[str] = None
    valid_hours: Optional[float] = None  # e.g. 2 = 2 hours; 24 = 1 day; None = no expiry


class ClassroomJoin(BaseModel):
    display_name: str = "Student"


class ClassroomStart(BaseModel):
    material_id: Optional[str] = None
    text: Optional[str] = None


def _gen_room_code() -> str:
    import string
    return "".join(random.choices(string.ascii_uppercase + string.digits, k=5))


def _classroom_expired(room: dict) -> bool:
    exp = room.get("expires_at")
    if not exp:
        return False
    exp_dt = _parse_dt(exp)
    return bool(exp_dt and exp_dt < datetime.now(timezone.utc))


def _assert_classroom_active(room: dict) -> None:
    if _classroom_expired(room):
        raise HTTPException(status_code=410, detail="Room code has expired — ask your teacher for a new code")


@api_router.get("/teacher/heatmap")
async def teacher_heatmap(user: Optional[User] = Depends(get_current_user)):
    """Class-wide memory strength grid for teacher dashboard."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    roster_resp = await teacher_roster(user)
    students = roster_resp.get("students", [])
    buckets = {"critical": 0, "weak": 0, "ok": 0, "strong": 0}
    cells = []
    for s in students:
        pct = s.get("readiness_percent", 0)
        if pct < 40:
            bucket = "critical"
        elif pct < 60:
            bucket = "weak"
        elif pct < 80:
            bucket = "ok"
        else:
            bucket = "strong"
        buckets[bucket] += 1
        cells.append({**s, "bucket": bucket})
    return {"buckets": buckets, "students": cells, "room_count": roster_resp.get("room_count", 0)}


@api_router.get("/teacher/roster")
async def teacher_roster(user: Optional[User] = Depends(get_current_user)):
    """Simple student readiness list across all teacher classrooms."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    rooms = await db.classrooms.find({"teacher_id": user.user_id}, {"_id": 0}).to_list(100)
    now = datetime.now(timezone.utc)
    roster = []
    seen = set()
    for room in rooms:
        for p in room.get("participants", []):
            guest_id = p.get("guest_id") or p.get("display_name", "")
            key = p.get("guest_id") or p.get("display_name")
            if not key or key in seen:
                continue
            seen.add(key)
            student_owner = f"guest_{guest_id}" if guest_id and not str(guest_id).startswith("guest_") else guest_id
            if not student_owner or student_owner == "guest_":
                student_owner = None
            readiness = 0.0
            units_due = 0
            track_count = 0
            if student_owner:
                tracks = await db.learning_tracks.find(_student_tracks_filter(student_owner), {"_id": 0}).to_list(50)
                track_count = len(tracks)
                if tracks:
                    tids = [t["track_id"] for t in tracks]
                    units = await db.knowledge_units.find({"track_id": {"$in": tids}}, {"_id": 0}).to_list(500)
                    strengths = [_unit_memory_strength(u, now) for u in units]
                    readiness = round(sum(strengths) / len(strengths), 1) if strengths else 0.0
                    units_due = sum(
                        1 for u in units
                        if u.get("review_count", 0) > 0 and _unit_memory_strength(u, now) < MEMORY_THRESHOLD
                    )
            roster.append({
                "display_name": p.get("display_name", "Student"),
                "guest_id": guest_id,
                "room_code": room.get("room_code"),
                "room_name": room.get("name"),
                "readiness_percent": readiness,
                "units_due": units_due,
                "track_count": track_count,
                "last_score": p.get("score"),
                "status": p.get("status", "waiting"),
            })
    roster.sort(key=lambda r: r.get("readiness_percent", 0))
    return {"students": roster, "room_count": len(rooms)}


@api_router.get("/teacher/calendar")
async def teacher_calendar(user: Optional[User] = Depends(get_current_user), days_ahead: int = 21):
    """Class-wide deadlines + review-load forecast for the Teacher dashboard
    calendar — aggregates due dates and daily review load across every
    student in every one of this teacher's classrooms."""
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    rooms = await db.classrooms.find({"teacher_id": user.user_id}, {"_id": 0}).to_list(100)
    now = datetime.now(timezone.utc)
    horizon = now + timedelta(days=days_ahead)

    student_owners = set()
    for room in rooms:
        for p in room.get("participants", []):
            guest_id = p.get("guest_id")
            if not guest_id:
                continue
            student_owners.add(f"guest_{guest_id}" if not str(guest_id).startswith("guest_") else guest_id)

    events = []
    all_units = []
    if student_owners:
        tracks = await db.learning_tracks.find(
            {"$or": [_student_tracks_filter(o) for o in student_owners]}, {"_id": 0}
        ).to_list(500)
        tids = [t["track_id"] for t in tracks]
        seen_tids = set()
        for t in sorted(tracks, key=lambda t: t.get("due_date") or ""):
            if t["track_id"] in seen_tids:
                continue
            seen_tids.add(t["track_id"])
            due_dt = _parse_dt(t.get("due_date")) if t.get("due_date") else None
            if due_dt and now <= due_dt <= horizon:
                events.append({
                    "date": due_dt.date().isoformat(),
                    "type": "deadline",
                    "track_id": t["track_id"],
                    "title": t.get("title", "Untitled"),
                    "track_type": t.get("track_type"),
                    "urgency_tier": compute_urgency_tier(t.get("track_type"), due_dt, now),
                })
        if tids:
            all_units = await db.knowledge_units.find({"track_id": {"$in": tids}}, {"_id": 0}).to_list(5000)

    daily_load = _bucket_units_by_due_date(all_units, days_ahead)

    return {
        "events": sorted(events, key=lambda e: e["date"]),
        "daily_load": [{"date": d, "units_due": c} for d, c in sorted(daily_load.items())],
        "student_count": len(student_owners),
        "generated_at": now.isoformat(),
    }


@api_router.post("/classrooms")
async def create_classroom(
    payload: ClassroomCreate,
    user: Optional[User] = Depends(get_current_user),
):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    if user.role == "teacher":
        plan = await get_teacher_plan(db, user.user_id)
        room_count = await db.classrooms.count_documents({"teacher_id": user.user_id})
        if room_count >= plan.get("classroom_limit", 1):
            raise HTTPException(status_code=402, detail="Classroom limit reached — upgrade your teacher plan")
    room_code = _gen_room_code()
    expires_at = None
    if payload.valid_hours and payload.valid_hours > 0:
        expires_at = (datetime.now(timezone.utc) + timedelta(hours=payload.valid_hours)).isoformat()
    room = Classroom(
        room_code=room_code,
        teacher_id=user.user_id,
        name=payload.name,
        mode=payload.mode,
        material_id=payload.material_id,
        expires_at=expires_at,
    )
    doc = room.model_dump()
    doc["created_at"] = doc["created_at"].isoformat()
    await db.classrooms.insert_one(doc)
    return room


@api_router.get("/classrooms")
async def list_classrooms(user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    docs = await db.classrooms.find({"teacher_id": user.user_id}, {"_id": 0}).sort("created_at", -1).to_list(100)
    for d in docs:
        if _classroom_expired(d):
            d["status"] = "expired"
    return docs


@api_router.get("/classrooms/{room_code}")
async def get_classroom(room_code: str):
    doc = await db.classrooms.find_one({"room_code": room_code.upper()}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Room not found")
    if _classroom_expired(doc):
        doc["status"] = "expired"
    return doc


@api_router.post("/classrooms/{room_code}/join")
async def join_classroom(
    room_code: str,
    payload: ClassroomJoin,
    owner: str = Depends(get_owner_id),
    x_guest_id: Optional[str] = Header(default=None, alias="X-Guest-Id"),
):
    code = room_code.upper()
    room = await db.classrooms.find_one({"room_code": code}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    _assert_classroom_active(room)
    if room.get("status") == "finished":
        raise HTTPException(status_code=400, detail="Raid already finished")
    guest_key = x_guest_id or owner.replace("guest_", "")
    participant = {
        "guest_id": guest_key,
        "owner": owner,
        "display_name": payload.display_name,
        "score": 0,
        "progress": 0,
    }
    participants = [p for p in room.get("participants", []) if p.get("guest_id") != guest_key]
    participants.append(participant)
    await db.classrooms.update_one({"room_code": code}, {"$set": {"participants": participants}})
    await register_teacher_seat(db, room["teacher_id"], code, owner)
    await classroom_manager.broadcast(code, {"type": "participant_joined", "participant": participant})
    return {"ok": True, "room_code": code, "status": room.get("status", "waiting")}


@api_router.post("/classrooms/{room_code}/start")
async def start_classroom_raid(
    room_code: str,
    payload: ClassroomStart,
    user: Optional[User] = Depends(get_current_user),
):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    code = room_code.upper()
    room = await db.classrooms.find_one({"room_code": code, "teacher_id": user.user_id}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    _assert_classroom_active(room)
    text = payload.text
    if payload.material_id:
        mat = await db.materials.find_one({"material_id": payload.material_id}, {"_id": 0})
        if mat:
            text = mat.get("text", "")
    elif room.get("material_id"):
        mat = await db.materials.find_one({"material_id": room["material_id"]}, {"_id": 0})
        if mat:
            text = mat.get("text", "")
    if not text or len(text.strip()) < 8:
        raise HTTPException(status_code=400, detail="Provide material or text for the raid")
    result = await generate_game(GenerateGameRequest(text=text, mode=room.get("mode", "quiz")))
    game = result["game"]
    game["boss_max_hp"] = 100
    await db.classrooms.update_one(
        {"room_code": code},
        {"$set": {"status": "active", "active_game": game, "started_at": datetime.now(timezone.utc).isoformat()}},
    )
    await classroom_manager.broadcast(code, {"type": "game_started", "game": game})
    return {"ok": True, "game": game}


@api_router.post("/classrooms/{room_code}/finish")
async def finish_classroom_raid(room_code: str, user: Optional[User] = Depends(get_current_user)):
    if not user:
        raise HTTPException(status_code=401, detail="Login required")
    code = room_code.upper()
    room = await db.classrooms.find_one({"room_code": code, "teacher_id": user.user_id}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    await db.classrooms.update_one(
        {"room_code": code},
        {"$set": {"status": "finished", "finished_at": datetime.now(timezone.utc).isoformat()}},
    )
    await classroom_manager.broadcast(code, {"type": "raid_finished"})
    return {"ok": True}


@api_router.post("/classrooms/{room_code}/progress")
async def update_raid_progress(
    room_code: str,
    payload: dict,
    owner: str = Depends(get_owner_id),
    x_guest_id: Optional[str] = Header(default=None, alias="X-Guest-Id"),
):
    code = room_code.upper()
    guest_key = x_guest_id or owner.replace("guest_", "")
    room = await db.classrooms.find_one({"room_code": code}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    participants = room.get("participants", [])
    for p in participants:
        if p.get("guest_id") == guest_key:
            p["score"] = payload.get("score", p.get("score", 0))
            p["progress"] = payload.get("progress", p.get("progress", 0))
    await db.classrooms.update_one({"room_code": code}, {"$set": {"participants": participants}})
    await classroom_manager.broadcast(code, {"type": "progress_update", "guest_id": guest_key, **payload})
    return {"ok": True}


# ============== ROOT ==============
@api_router.get("/")
async def root():
    dev_on = os.environ.get("DEV_AUTH_ENABLED", "").lower() == "true"
    return {
        "message": "AI Cognitive Adventure Platform API",
        "ok": True,
        "dev_auth_enabled": dev_on,
    }


@app.on_event("startup")
async def _warm_mongo():
    """Atlas cold connect can take ~2s; warm the pool so /health stays fast."""
    try:
        await asyncio.wait_for(db.command("ping"), timeout=10.0)
    except Exception:
        pass


@api_router.get("/health")
async def health():
    """Quick readiness check — frontend can surface a useful error when Mongo is down."""
    mongo_ok = False
    try:
        # Atlas first ping can exceed 2s; 5s keeps local dev responsive without false negatives.
        await asyncio.wait_for(db.command("ping"), timeout=5.0)
        mongo_ok = True
    except Exception:
        pass
    return {
        "ok": mongo_ok,
        "mongo": mongo_ok,
        "gemini": GEMINI_ENABLED,
        "vision": VISION_ENABLED,
        "resend": bool(os.environ.get("RESEND_API_KEY", "")) and not str(os.environ.get("RESEND_API_KEY", "")).startswith("PASTE"),
        "firebase": bool(os.environ.get("FIREBASE_SERVICE_ACCOUNT_JSON", "")),
        "message": "ready" if mongo_ok else "MongoDB unreachable — start mongod or check MONGO_URL",
    }


register_teacher_billing_routes(api_router, db, get_current_user)
register_iap_routes(api_router, db, get_current_user, get_owner_id)

app.include_router(api_router)


@app.websocket("/api/ws/classrooms/{room_code}")
async def classroom_websocket(websocket: WebSocket, room_code: str):
    code = room_code.upper()
    await classroom_manager.connect(code, websocket)
    try:
        while True:
            data = await websocket.receive_json()
            if data.get("type") == "ping":
                await websocket.send_json({"type": "pong"})
    except WebSocketDisconnect:
        classroom_manager.disconnect(code, websocket)


def _cors_origins() -> list[str]:
    raw = os.environ.get("CORS_ORIGINS", "*").strip()
    origins: list[str] = []
    if raw in ("", "*"):
        origins = ["http://localhost:3000", "http://127.0.0.1:3000"]
    else:
        origins = [o.strip() for o in raw.split(",") if o.strip()]
    app_base = os.environ.get("APP_BASE_URL", "").strip().rstrip("/")
    if app_base and app_base not in origins:
        origins.append(app_base)
    return origins


app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=_cors_origins(),
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
